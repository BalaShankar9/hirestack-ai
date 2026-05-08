"""
S15-P1: smoke tests for the PPT orchestrator.

Strategy: drive everything with a stub AIClient so tests are offline-only
and deterministic. Validate that:
- OutlinePlanner coerces an LLM payload into a valid DeckSpec.
- SlideComposer produces a valid .pptx (parsed back via python-pptx).
- PPTOrchestrator.generate end-to-end returns reasonable output.
- Planner falls back to deterministic stub when LLM raises.
"""
from __future__ import annotations

import io
from typing import Any, Dict, Optional

import pytest

pytestmark = pytest.mark.asyncio

# Skip the whole module if python-pptx isn't available (CI without optional deps).
pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

from ai_engine.agents.ppt import (  # noqa: E402
    DeckSpec,
    OutlinePlanner,
    PPTOrchestrator,
    SlideComposer,
    SlideKind,
    SlideSpec,
)
from ai_engine.agents.orchestration import (  # noqa: E402
    ORCHESTRATION_PROGRESS_SCHEMA_VERSION,
    PPT_GENERATION_PHASE_ORDER,
)


class _StubClient:
    """Minimal AIClient surrogate exposing only complete_json."""

    def __init__(self, payload: Dict[str, Any]) -> None:
        self.payload = payload
        self.calls = 0

    async def complete_json(self, **kwargs):
        self.calls += 1
        return self.payload


class _RaisingClient:
    async def complete_json(self, **kwargs):
        raise RuntimeError("LLM unavailable for test")


def _good_payload() -> Dict[str, Any]:
    return {
        "title": "HireStack AI Investor Pitch",
        "subtitle": "Series A — 2026",
        "theme": "modern",
        "accent_color": "#2563EB",
        "slides": [
            {"kind": "title", "title": "HireStack AI", "subtitle": "Investor Pitch", "notes": "Hello"},
            {"kind": "section", "title": "The Problem"},
            {"kind": "content", "title": "Hiring is broken",
             "bullets": ["Slow", "Biased", "Expensive"], "notes": "Talk about pain"},
            {"kind": "two_column", "title": "Today vs Tomorrow",
             "bullets": ["Manual", "Slow"], "bullets_right": ["AI", "Fast"]},
            {"kind": "quote", "body": "AI eats hiring", "attribution": "VC Friend"},
            {"kind": "chart", "title": "Growth", "caption": "MRR by quarter",
             "chart": {"kind": "line", "title": "MRR",
                       "series": [{"name": "MRR", "data": [10, 20, 35, 60, 100]}],
                       "categories": ["Q1", "Q2", "Q3", "Q4", "Q5"]}},
            {"kind": "image", "title": "Our Team",
             "image": {"query": "diverse engineering team", "alt_text": "team"}},
            {"kind": "image_text", "title": "Product",
             "bullets": ["Resume tailoring", "ATS scan"],
             "image": {"query": "modern dashboard"}},
            {"kind": "closing", "title": "Thank you", "body": "Questions?"},
        ],
    }


# ────────────────────────────────────────────────────────────────────
#  OutlinePlanner
# ────────────────────────────────────────────────────────────────────

async def test_planner_coerces_payload_into_deckspec():
    planner = OutlinePlanner(ai_client=_StubClient(_good_payload()))
    deck = await planner.plan(topic="HireStack AI", audience="investors", slide_count=9)
    assert isinstance(deck, DeckSpec)
    assert deck.title == "HireStack AI Investor Pitch"
    assert deck.slide_count == 9
    assert deck.slides[0].kind == SlideKind.title
    assert deck.slides[-1].kind == SlideKind.closing
    chart_slide = next(s for s in deck.slides if s.kind == SlideKind.chart)
    assert chart_slide.chart is not None
    assert chart_slide.chart.kind == "line"
    assert chart_slide.chart.series[0]["data"] == [10, 20, 35, 60, 100]


async def test_planner_falls_back_to_stub_when_llm_raises():
    planner = OutlinePlanner(ai_client=_RaisingClient())
    deck = await planner.plan(topic="Quantum Tea", audience="execs", slide_count=5)
    assert deck.slide_count == 5
    assert deck.slides[0].kind == SlideKind.title
    assert deck.slides[-1].kind == SlideKind.closing
    assert deck.title == "Quantum Tea"


async def test_planner_clamps_slide_count_bounds():
    planner = OutlinePlanner(ai_client=_RaisingClient())
    too_small = await planner.plan(topic="X", slide_count=1)
    too_big = await planner.plan(topic="X", slide_count=999)
    assert too_small.slide_count >= 3
    assert too_big.slide_count <= 30


async def test_planner_drops_unknown_slide_kind_to_content():
    payload = {
        "title": "Mixed",
        "slides": [
            {"kind": "garbage", "title": "Should become content", "bullets": ["a", "b"]},
        ],
    }
    planner = OutlinePlanner(ai_client=_StubClient(payload))
    deck = await planner.plan(topic="X", slide_count=3)
    assert deck.slides[0].kind == SlideKind.content


# ────────────────────────────────────────────────────────────────────
#  SlideComposer
# ────────────────────────────────────────────────────────────────────

async def test_composer_produces_valid_pptx_bytes():
    deck = DeckSpec(
        title="Smoke",
        slides=[
            SlideSpec(kind=SlideKind.title, title="Smoke Deck", subtitle="MVP"),
            SlideSpec(kind=SlideKind.content, title="Bullets",
                      bullets=["one", "two", "three"], notes="Speak slowly"),
            SlideSpec(kind=SlideKind.closing, title="Thanks", body="Q?"),
        ],
    )
    composer = SlideComposer()
    pptx_bytes = await composer.compose(deck)
    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 5_000  # any real .pptx is comfortably bigger
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 3


async def test_composer_handles_all_slide_kinds_without_errors():
    deck = DeckSpec(
        title="All kinds",
        slides=[
            SlideSpec(kind=SlideKind.title, title="T"),
            SlideSpec(kind=SlideKind.section, title="S"),
            SlideSpec(kind=SlideKind.content, title="C", bullets=["b1", "b2"]),
            SlideSpec(kind=SlideKind.two_column, title="2C",
                      bullets=["L1"], bullets_right=["R1"]),
            SlideSpec(kind=SlideKind.quote, body="Be bold", attribution="Anon"),
            SlideSpec(kind=SlideKind.chart, title="Chart"),  # no chart obj → placeholder
            SlideSpec(kind=SlideKind.image, title="Img"),    # no image → placeholder
            SlideSpec(kind=SlideKind.image_text, title="IT", bullets=["x"]),
            SlideSpec(kind=SlideKind.closing, title="Bye"),
        ],
    )
    composer = SlideComposer()
    pptx_bytes = await composer.compose(deck)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 9


async def test_composer_writes_speaker_notes():
    deck = DeckSpec(
        title="Notes",
        slides=[SlideSpec(kind=SlideKind.title, title="T", notes="hello speaker")],
    )
    pptx_bytes = await SlideComposer().compose(deck)
    prs = Presentation(io.BytesIO(pptx_bytes))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "hello speaker" in notes_text


async def test_composer_accepts_accent_color_override():
    deck = DeckSpec(
        title="Accent",
        accent_color="#FF00AA",
        slides=[SlideSpec(kind=SlideKind.title, title="T")],
    )
    # Should not raise even on weird accent.
    pptx_bytes = await SlideComposer().compose(deck)
    assert len(pptx_bytes) > 5_000


# ────────────────────────────────────────────────────────────────────
#  Orchestrator end-to-end
# ────────────────────────────────────────────────────────────────────

async def test_orchestrator_end_to_end_returns_pptresult():
    orch = PPTOrchestrator(ai_client=_StubClient(_good_payload()))
    result = await orch.generate(topic="Pitch", audience="vcs", slide_count=9)
    assert result.size_bytes > 5_000
    assert result.slide_count == 9
    assert result.latency_ms >= 0
    assert tuple(result.phase_latencies.keys())[0] == PPT_GENERATION_PHASE_ORDER[0]
    assert result.phase_statuses["outline"] == "completed"
    assert result.phase_statuses["composition"] == "completed"
    prs = Presentation(io.BytesIO(result.pptx_bytes))
    assert len(prs.slides) == 9


async def test_orchestrator_rejects_blank_topic():
    orch = PPTOrchestrator(ai_client=_StubClient(_good_payload()))
    with pytest.raises(ValueError):
        await orch.generate(topic="   ", slide_count=5)


async def test_orchestrator_generate_from_deck_skips_planner():
    deck = DeckSpec(title="X", slides=[SlideSpec(kind=SlideKind.title, title="T")])
    stub = _StubClient(_good_payload())
    orch = PPTOrchestrator(ai_client=stub)
    result = await orch.generate_from_deck(deck)
    assert stub.calls == 0  # planner was bypassed
    assert result.size_bytes > 1_000
    assert "outline" not in result.phase_statuses


async def test_orchestrator_progress_callback_uses_canonical_phase_names():
    seen = []
    payloads = []

    def _capture(progress):
        if progress.phase is not None:
            seen.append(progress.phase)
            payloads.append(progress.to_payload())

    orch = PPTOrchestrator(
        ai_client=_StubClient(_good_payload()),
        progress_callback=_capture,
    )

    await orch.generate(topic="Pitch", audience="vcs", slide_count=9)

    assert seen
    assert seen[0] == PPT_GENERATION_PHASE_ORDER[0]
    assert all(phase in PPT_GENERATION_PHASE_ORDER for phase in seen)
    assert payloads
    assert all(payload["pipeline_name"] == "ppt_generation" for payload in payloads)
    assert all(payload["event_type"] == "progress" for payload in payloads)
    assert all(payload["schema_version"] == ORCHESTRATION_PROGRESS_SCHEMA_VERSION for payload in payloads)
    assert all(payload["phase"] == payload["stage"] for payload in payloads)
