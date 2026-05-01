"""
S15-P2: smoke tests for the matplotlib chart galaxy.

Each test renders a representative ChartSpec for one chart family and
asserts the renderer produced a non-trivial PNG byte buffer. We don't
do pixel-level snapshots (matplotlib output is platform-sensitive); we
just verify (a) bytes returned, (b) PNG signature, (c) the orchestrator
composes a deck containing a real chart image rather than a placeholder.
"""
from __future__ import annotations

import io

import pytest

pytest.importorskip("pptx")
pytest.importorskip("matplotlib")
from pptx import Presentation  # noqa: E402

from ai_engine.agents.ppt import (  # noqa: E402
    ChartRenderer,
    ChartSelector,
    DeckSpec,
    PPTOrchestrator,
    SlideComposer,
    SlideKind,
    SlideSpec,
)
from ai_engine.agents.ppt.schemas import ChartKind, ChartSpec  # noqa: E402


PNG_SIG = b"\x89PNG\r\n\x1a\n"

# Async-mark the per-renderer + e2e tests; the ChartSelector tests below are sync.
async_test = pytest.mark.asyncio


def _ok_png(b):
    assert b is not None and isinstance(b, (bytes, bytearray))
    assert b[:8] == PNG_SIG
    assert len(b) > 1500  # any chart has >1KB of pixels


# ─── per-kind smoke renders ──────────────────────────────────────────

@async_test
async def test_render_line():
    spec = ChartSpec(kind="line", title="MRR",
                     series=[{"name": "MRR", "data": [10, 20, 35, 60, 100]}],
                     categories=["Q1", "Q2", "Q3", "Q4", "Q5"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_area_multi_series():
    spec = ChartSpec(kind="area",
                     series=[{"name": "A", "data": [1, 2, 3, 4]},
                             {"name": "B", "data": [4, 3, 2, 1]}])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_bar_horizontal():
    spec = ChartSpec(kind="bar", title="Top markets",
                     series=[{"name": "Rev", "data": [120, 90, 70, 40]}],
                     categories=["US", "EU", "APAC", "LATAM"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_column_grouped():
    spec = ChartSpec(kind="column",
                     series=[{"name": "2025", "data": [10, 20, 30]},
                             {"name": "2026", "data": [15, 25, 35]}],
                     categories=["Q1", "Q2", "Q3"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_stacked_bar():
    spec = ChartSpec(kind="stacked_bar",
                     series=[{"name": "Free", "data": [10, 20, 30]},
                             {"name": "Pro", "data": [5, 15, 25]}],
                     categories=["Q1", "Q2", "Q3"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_scatter_pairs():
    spec = ChartSpec(kind="scatter",
                     series=[{"name": "S", "data": [[1, 2], [2, 4], [3, 5], [4, 7]]}])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_bubble_triples():
    spec = ChartSpec(kind="bubble",
                     series=[{"name": "B", "data": [[1, 2, 5], [2, 3, 12], [3, 4, 20]]}])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_pie():
    spec = ChartSpec(kind="pie",
                     series=[{"name": "Share", "data": [40, 30, 20, 10]}],
                     categories=["A", "B", "C", "D"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_donut():
    spec = ChartSpec(kind="donut",
                     series=[{"name": "Share", "data": [55, 25, 20]}],
                     categories=["X", "Y", "Z"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_histogram():
    spec = ChartSpec(kind="histogram",
                     series=[{"name": "Latency", "data": [10, 11, 12, 12, 13, 13, 14, 15, 18, 22, 25, 30]}])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_box_multi():
    spec = ChartSpec(kind="box",
                     series=[{"name": "A", "data": [1, 2, 3, 4, 5, 6, 7]},
                             {"name": "B", "data": [2, 4, 6, 8, 10, 12, 14]}])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_heatmap_matrix():
    spec = ChartSpec(kind="heatmap",
                     matrix=[[1, 2, 3], [4, 5, 6], [7, 8, 9]],
                     categories=["c1", "c2", "c3"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_waterfall():
    spec = ChartSpec(kind="waterfall",
                     series=[{"name": "Δ", "data": [10, 5, -3, 8, -2, 12]}],
                     categories=["Start", "Sales", "Refund", "Upsell", "Churn", "End"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_radar():
    spec = ChartSpec(kind="radar",
                     series=[{"name": "Profile", "data": [4, 5, 3, 4, 5, 2]}],
                     categories=["Speed", "Cost", "UX", "Support", "Trust", "Reach"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_funnel():
    spec = ChartSpec(kind="funnel",
                     series=[{"name": "Pipeline", "data": [1000, 500, 200, 80, 30]}],
                     categories=["Visit", "Signup", "Trial", "Demo", "Close"])
    _ok_png(await ChartRenderer().render(spec))


@async_test
async def test_render_unknown_kind_returns_none():
    spec = ChartSpec(kind="totallyfake", series=[{"name": "x", "data": [1, 2]}])
    assert await ChartRenderer().render(spec) is None


# ─── ChartSelector heuristics ────────────────────────────────────────

def test_selector_pie_for_share_summing_100():
    pick = ChartSelector.suggest(
        series=[{"name": "Share", "data": [50, 30, 20]}],
        categories=["A", "B", "C"],
    )
    assert pick == ChartKind.pie


def test_selector_line_for_time_series():
    pick = ChartSelector.suggest(
        series=[{"name": "Rev", "data": [10, 20, 30]}],
        categories=["Q1 2025", "Q2 2025", "Q3 2025"],
    )
    assert pick == ChartKind.line


def test_selector_column_for_simple_categorical():
    pick = ChartSelector.suggest(
        series=[{"name": "Rev", "data": [10, 20, 30, 40, 50, 60, 70]}],
        categories=["A", "B", "C", "D", "E", "F", "G"],
    )
    assert pick == ChartKind.column


def test_selector_scatter_for_xy_pairs():
    pick = ChartSelector.suggest(
        series=[{"name": "S", "data": [[1, 2], [3, 4], [5, 6]]}],
    )
    assert pick == ChartKind.scatter


def test_selector_bubble_for_xyz_triples():
    pick = ChartSelector.suggest(
        series=[{"name": "B", "data": [[1, 2, 3], [2, 3, 5]]}],
    )
    assert pick == ChartKind.bubble


def test_selector_funnel_via_hint():
    pick = ChartSelector.suggest(
        series=[{"name": "F", "data": [1000, 500, 100]}],
        hint="conversion",
    )
    assert pick == ChartKind.funnel


# ─── End-to-end: chart shows up inside the .pptx ─────────────────────

@async_test
async def test_composer_embeds_real_chart_when_renderer_provided():
    deck = DeckSpec(
        title="Real chart",
        slides=[
            SlideSpec(
                kind=SlideKind.chart,
                title="Q growth",
                chart=ChartSpec(
                    kind="column", title="Q growth",
                    series=[{"name": "Rev", "data": [10, 20, 30]}],
                    categories=["Q1", "Q2", "Q3"],
                ),
            ),
        ],
    )
    composer = SlideComposer(chart_renderer=ChartRenderer())
    pptx_bytes = await composer.compose(deck)
    prs = Presentation(io.BytesIO(pptx_bytes))
    # At least one picture shape should now exist on the chart slide.
    pic_count = sum(1 for shp in prs.slides[0].shapes if shp.shape_type is not None
                    and getattr(shp, "image", None) is not None)
    assert pic_count >= 1


@async_test
async def test_orchestrator_default_renderer_is_chart_galaxy():
    orch = PPTOrchestrator(ai_client=None)
    # composer should have inherited a non-None chart_renderer
    assert orch.composer.chart_renderer is not None
