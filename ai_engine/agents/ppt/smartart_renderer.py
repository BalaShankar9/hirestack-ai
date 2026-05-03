"""
SmartArtRenderer — Process flows, hierarchies, cycles, timelines.

Provides diagram types common in business presentations:
- Process: Linear step-by-step flow
- Cycle: Circular/recurring processes
- Hierarchy: Org charts, decision trees
- Pyramid: Layered structures (Maslow, food pyramid)
- Timeline: Gantt-style or milestone views
- Funnel: Conversion flows

All diagrams render as native PPTX shapes (not images) for editability.

Public API:
    SmartArtRenderer.render(slide, kind, data, theme) -> None
"""
from __future__ import annotations

import logging
from typing import Any, Dict, List, Optional, Tuple

from ai_engine.agents.ppt.schemas import SlideSpec

logger = logging.getLogger(__name__)


def _hex_to_rgb(value: str) -> Tuple[int, int, int]:
    """Convert #RRGGBB to RGB tuple."""
    h = value.strip().lstrip("#")
    if len(h) != 6:
        return (37, 99, 235)  # Default blue
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class SmartArtRenderer:
    """Render SmartArt-style diagrams using native PPTX shapes."""

    # Diagram type configurations
    DIAGRAM_TYPES = {
        "process", "cycle", "hierarchy", "pyramid", "timeline", "funnel"
    }

    def __init__(self) -> None:
        pass

    def can_render(self, diagram_type: str) -> bool:
        """Check if we can render this diagram type."""
        return diagram_type.lower().strip() in self.DIAGRAM_TYPES

    def render(
        self,
        slide,
        diagram_type: str,
        items: List[Dict[str, Any]],
        *,
        accent_hex: str = "#2563EB",
        title_c: str = "#0F172A",
        body_c: str = "#334155",
    ) -> bool:
        """
        Render a diagram to the slide.

        Args:
            slide: PPTX slide object
            diagram_type: One of process, cycle, hierarchy, pyramid, timeline, funnel
            items: List of {"label": str, "description": str, "level": int}
            accent_hex: Theme accent color
            title_c: Title text color
            body_c: Body text color

        Returns:
            True if rendered successfully
        """
        kind = diagram_type.lower().strip()

        try:
            if kind == "process":
                return self._render_process(slide, items, accent_hex, title_c, body_c)
            elif kind == "cycle":
                return self._render_cycle(slide, items, accent_hex, title_c, body_c)
            elif kind == "hierarchy":
                return self._render_hierarchy(slide, items, accent_hex, title_c, body_c)
            elif kind == "pyramid":
                return self._render_pyramid(slide, items, accent_hex, title_c, body_c)
            elif kind == "timeline":
                return self._render_timeline(slide, items, accent_hex, title_c, body_c)
            elif kind == "funnel":
                return self._render_funnel(slide, items, accent_hex, title_c, body_c)
            else:
                logger.debug("smartart_unknown_type: %s", kind)
                return False
        except Exception as exc:
            logger.warning("smartart_render_failed: %s", str(exc)[:200])
            return False

    def _render_process(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render horizontal process flow with arrows."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not items:
            return False

        n = min(len(items), 6)  # Max 6 steps
        box_width = Inches(10.0 / n)
        box_height = Inches(1.2)
        start_left = Inches(1.5)
        top = Inches(2.5)
        accent_rgb = _hex_to_rgb(accent_hex)

        for i, item in enumerate(items[:n]):
            left = start_left + (box_width + Inches(0.3)) * i

            # Draw box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width - Inches(0.2), box_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*accent_rgb if i == 0 else (241, 245, 249))
            shape.line.color.rgb = RGBColor(*accent_rgb)
            shape.line.width = Pt(2)

            # Add text
            tf = shape.text_frame
            tf.text = item.get("label", f"Step {i+1}")
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*(_hex_to_rgb(title_c) if i == 0 else _hex_to_rgb(body_c)))
            p.alignment = 1  # Center

            # Add arrow (except for last item)
            if i < n - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    left + box_width - Inches(0.1), top + Inches(0.4),
                    Inches(0.3), Inches(0.4),
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(*accent_rgb)
                arrow.line.fill.background()

        return True

    def _render_cycle(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render circular cycle diagram."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import math

        if not items or len(items) < 3:
            return False

        n = min(len(items), 6)
        center_x, center_y = Inches(6.67), Inches(3.75)  # Slide center
        radius = Inches(2.5)
        box_size = Inches(1.5)
        accent_rgb = _hex_to_rgb(accent_hex)

        # Color palette
        colors = [accent_rgb]
        colors.extend([
            (59, 130, 246),   # Blue
            (16, 185, 129),   # Green
            (245, 158, 11),   # Amber
            (139, 92, 246),   # Purple
            (236, 72, 153),   # Pink
        ])

        for i, item in enumerate(items[:n]):
            angle = (2 * math.pi * i / n) - math.pi / 2  # Start at top
            left = center_x + radius * math.cos(angle) - box_size / 2
            top = center_y + radius * math.sin(angle) - box_size / 2

            # Draw circle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, box_size, box_size,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*colors[i])
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.text = item.get("label", f"{i+1}")
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = 1  # Center

        # Draw center hub
        hub = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            center_x - Inches(0.6), center_y - Inches(0.6),
            Inches(1.2), Inches(1.2),
        )
        hub.fill.solid()
        hub.fill.fore_color.rgb = RGBColor(*accent_rgb)
        hub.line.fill.background()

        return True

    def _render_hierarchy(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render org chart / tree hierarchy."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not items:
            return False

        # Sort by level
        levels: Dict[int, List[Dict]] = {}
        for item in items:
            level = item.get("level", 0)
            levels.setdefault(level, []).append(item)

        accent_rgb = _hex_to_rgb(accent_hex)
        box_width, box_height = Inches(2.0), Inches(0.6)
        start_y = Inches(1.5)
        level_height = Inches(1.2)

        max_level = max(levels.keys()) if levels else 0
        for level in range(max_level + 1):
            level_items = levels.get(level, [])
            if not level_items:
                continue

            n = len(level_items)
            total_width = n * box_width + (n - 1) * Inches(0.5)
            start_x = (Inches(13.333) - total_width) / 2

            for i, item in enumerate(level_items):
                left = start_x + i * (box_width + Inches(0.5))
                top = start_y + level * level_height

                # Draw box
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, box_width, box_height,
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(241, 245, 249) if level > 0 else RGBColor(*accent_rgb)
                shape.line.color.rgb = RGBColor(*accent_rgb)

                # Text
                tf = shape.text_frame
                tf.text = item.get("label", f"Level {level}")
                p = tf.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(*(_hex_to_rgb(body_c) if level > 0 else (255, 255, 255)))
                p.alignment = 1

        return True

    def _render_pyramid(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render pyramid with stacked levels."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not items:
            return False

        n = min(len(items), 5)
        center_x = Inches(6.67)
        max_width = Inches(10)
        level_height = Inches(0.8)
        start_y = Inches(5.0)
        accent_rgb = _hex_to_rgb(accent_hex)

        for i in range(n):
            # Wider at bottom, narrower at top
            level_width = max_width * (n - i) / n
            left = center_x - level_width / 2
            top = start_y - (i + 1) * level_height

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, level_width, level_height,
            )
            shape.fill.solid()
            # Darker at bottom, lighter at top
            intensity = max(0.3, 1.0 - (i * 0.15))
            fill_rgb = tuple(int(c * intensity + 255 * (1 - intensity)) for c in accent_rgb)
            shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
            shape.line.color.rgb = RGBColor(*accent_rgb)

            # Text
            tf = shape.text_frame
            tf.text = items[n - 1 - i].get("label", f"Level {n-i}")
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*_hex_to_rgb(title_c))
            p.alignment = 1

        return True

    def _render_timeline(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render horizontal timeline with milestones."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not items:
            return False

        n = min(len(items), 6)
        start_x = Inches(1.5)
        end_x = Inches(12)
        center_y = Inches(4.0)
        accent_rgb = _hex_to_rgb(accent_hex)

        # Draw main line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            start_x, center_y - Inches(0.02),
            end_x - start_x, Inches(0.04),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*accent_rgb)
        line.line.fill.background()

        # Draw milestones
        spacing = (end_x - start_x) / (n - 1) if n > 1 else Inches(0)

        for i, item in enumerate(items[:n]):
            x = start_x + spacing * i

            # Draw milestone circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.15), center_y - Inches(0.15),
                Inches(0.3), Inches(0.3),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*accent_rgb)
            circle.line.color.rgb = RGBColor(255, 255, 255)
            circle.line.width = Pt(3)

            # Add label below
            label = slide.shapes.add_textbox(
                x - Inches(1), center_y + Inches(0.4),
                Inches(2), Inches(0.8),
            )
            tf = label.text_frame
            tf.text = item.get("label", f"M{i+1}")
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(*_hex_to_rgb(body_c))
            p.alignment = 1  # Center

        return True

    def _render_funnel(
        self, slide, items: List[Dict],
        accent_hex: str, title_c: str, body_c: str
    ) -> bool:
        """Render funnel with decreasing width."""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        if not items:
            return False

        n = min(len(items), 5)
        center_x = Inches(6.67)
        max_width = Inches(6)
        level_height = Inches(0.8)
        start_y = Inches(2.0)
        accent_rgb = _hex_to_rgb(accent_hex)

        for i in range(n):
            # Decreasing width
            width = max_width * (n - i) / n
            left = center_x - width / 2
            top = start_y + i * level_height

            # Use trapezoid if available, else rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, level_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
            shape.fill.fore_color.brightness = 1.0 - (i * 0.15)  # Lighter down the funnel
            shape.line.color.rgb = RGBColor(*accent_rgb)

            # Add text centered
            tf = shape.text_frame
            tf.text = items[i].get("label", f"Stage {i+1}")
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = 1

        return True


# Convenience function
def render_smartart(
    slide,
    diagram_type: str,
    items: List[Dict[str, Any]],
    *,
    accent_hex: str = "#2563EB",
    title_c: str = "#0F172A",
    body_c: str = "#334155",
) -> bool:
    """One-shot SmartArt render."""
    renderer = SmartArtRenderer()
    return renderer.render(slide, diagram_type, items, accent_hex=accent_hex, title_c=title_c, body_c=body_c)
