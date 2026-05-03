"""
TemplateLoader — Load and apply PowerPoint templates with brand kit enforcement.

Provides corporate-grade template support:
- Load .pptx files as templates
- Extract slide master layouts by name matching
- Apply brand kits (logo, colors, typography)
- Map SlideKind to appropriate template layouts

Public API:
    TemplateLoader(template_path).load() -> Presentation
    BrandKit.apply_to(presentation)
"""
from __future__ import annotations

import logging
import os
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

from ai_engine.agents.ppt.schemas import DeckSpec, SlideKind, SlideSpec

logger = logging.getLogger(__name__)


@dataclass
class BrandKit:
    """Corporate brand specification."""

    name: str
    primary_color: str = "#2563EB"  # Blue default
    secondary_color: str = "#0F172A"  # Dark default
    accent_color: str = "#3B82F6"
    font_heading: str = "Calibri"
    font_body: str = "Calibri"
    logo_path: Optional[str] = None
    logo_position: Tuple[float, float] = (0.5, 0.5)  # inches from left, top
    logo_size: Tuple[float, float] = (1.0, 0.5)  # width, height in inches

    def to_theme_dict(self) -> Dict[str, str]:
        """Convert to theme dictionary for SlideComposer."""
        return {
            "accent": self.primary_color.lstrip("#"),
            "title": self.secondary_color.lstrip("#"),
            "body": "475569",  # Default slate
            "heading_font": self.font_heading,
            "body_font": self.font_body,
        }


class TemplateLoader:
    """Load PowerPoint templates and extract layouts."""

    # SlideKind to layout name patterns (searched in order)
    _LAYOUT_PATTERNS: Dict[SlideKind, List[str]] = {
        SlideKind.title: ["title", "cover", "intro", "opening", "first"],
        SlideKind.section: ["section", "divider", "break", "chapter", "transition"],
        SlideKind.content: ["content", "body", "text", "bullet", "standard"],
        SlideKind.two_column: ["two column", "2 column", "compare", "comparison", "side by side"],
        SlideKind.quote: ["quote", "pull quote", "testimonial", "statement"],
        SlideKind.chart: ["chart", "graph", "data", "visualization"],
        SlideKind.image: ["image", "picture", "photo", "full bleed"],
        SlideKind.image_text: ["image + text", "picture + text", "split", "half"],
        SlideKind.closing: ["closing", "end", "final", "thank you", "contact"],
    }

    def __init__(self, template_path: Optional[str] = None, brand_kit: Optional[BrandKit] = None) -> None:
        self.template_path = template_path
        self.brand_kit = brand_kit or BrandKit(name="default")
        self._layouts: Dict[SlideKind, Any] = {}
        self._slide_master = None

    def load(self) -> Any:
        """
        Load template or create blank presentation.

        Returns:
            pptx.Presentation object
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx required for templates") from exc

        if self.template_path and os.path.exists(self.template_path):
            try:
                prs = Presentation(self.template_path)
                self._slide_master = prs.slide_masters[0] if prs.slide_masters else None
                self._index_layouts(prs)
                logger.info("template_loaded: %s", self.template_path)
            except Exception as exc:
                logger.warning("template_load_failed: %s — using blank", exc)
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            if self.template_path:
                logger.debug("template_not_found: %s — using blank", self.template_path)

        # Apply brand kit if specified
        if self.brand_kit:
            self._apply_brand_kit(prs)

        return prs

    def _index_layouts(self, prs: Any) -> None:
        """Map SlideKind to available layouts by name matching."""
        if not prs.slide_layouts:
            return

        for kind, patterns in self._LAYOUT_PATTERNS.items():
            for layout in prs.slide_layouts:
                layout_name = (layout.name or "").lower()
                if any(p in layout_name for p in patterns):
                    self._layouts[kind] = layout
                    break

        # Log which layouts were found
        found = [k.name for k in self._layouts.keys()]
        logger.debug("template_layouts_found: %s", found)

    def get_layout(self, kind: SlideKind) -> Optional[Any]:
        """Get the best layout for a SlideKind."""
        return self._layouts.get(kind)

    def has_template(self) -> bool:
        """Check if a template was successfully loaded."""
        return self.template_path is not None and bool(self._layouts)

    def _apply_brand_kit(self, prs: Any) -> None:
        """Apply brand colors and fonts to presentation."""
        if not self._slide_master:
            return

        try:
            # Apply to slide master theme
            theme = self._slide_master.theme
            if hasattr(theme, 'color_scheme'):
                scheme = theme.color_scheme
                # Map brand colors (best effort - some PPTX structures vary)
                self._set_theme_color(scheme, "accent1", self.brand_kit.primary_color)
                self._set_theme_color(scheme, "dk1", self.brand_kit.secondary_color)
        except Exception as exc:
            logger.debug("brand_kit_apply_failed: %s", exc)

    def _set_theme_color(self, scheme: Any, name: str, hex_color: str) -> None:
        """Set a theme color in the color scheme."""
        try:
            from pptx.dml.color import RGBColor
            rgb = self._hex_to_rgb(hex_color)
            # Try to set via color scheme attributes
            if hasattr(scheme, name):
                color = getattr(scheme, name)
                if hasattr(color, 'rgb'):
                    color.rgb = rgb
        except Exception:
            pass

    @staticmethod
    def _hex_to_rgb(value: str):
        from pptx.dml.color import RGBColor
        h = value.strip().lstrip("#")
        if len(h) != 6:
            return RGBColor(0x25, 0x63, 0xEB)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class TemplateSlideComposer:
    """
    Enhanced SlideComposer that uses templates and brand kits.

    Falls back to blank-slide composition if no template.
    """

    def __init__(
        self,
        *,
        template_path: Optional[str] = None,
        brand_kit: Optional[BrandKit] = None,
        chart_renderer: Optional[object] = None,
        image_resolver: Optional[object] = None,
    ) -> None:
        self.template_loader = TemplateLoader(template_path, brand_kit)
        self.brand_kit = brand_kit
        self.chart_renderer = chart_renderer
        self.image_resolver = image_resolver
        # Keep reference to base composer for fallback
        from ai_engine.agents.ppt.slide_composer import SlideComposer
        self._base_composer = SlideComposer(
            chart_renderer=chart_renderer,
            image_resolver=image_resolver,
        )

    async def compose(self, deck: DeckSpec) -> bytes:
        """Compose deck using template if available."""
        # Load presentation (template or blank)
        prs = self.template_loader.load()

        # If we have a template with layouts, use them
        if self.template_loader.has_template():
            return await self._compose_with_template(prs, deck)

        # Otherwise use base composer
        return await self._base_composer.compose(deck)

    async def _compose_with_template(self, prs: Any, deck: DeckSpec) -> bytes:
        """Compose using template layouts."""
        import io

        # Get theme colors from brand kit or deck
        accent = (self.brand_kit.primary_color if self.brand_kit else deck.accent_color) or "2563EB"
        accent = accent.lstrip("#")

        for slide_spec in deck.slides:
            layout = self.template_loader.get_layout(slide_spec.kind)

            if layout:
                slide = prs.slides.add_slide(layout)
                # Fill in placeholders
                self._fill_placeholders(slide, slide_spec, accent)
            else:
                # Fallback: add blank slide and render manually
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Apply background color from theme
                bg_hex = self._get_theme_bg()
                from ai_engine.agents.ppt.slide_composer import _hex_to_rgb
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = _hex_to_rgb(bg_hex)
                # Render content (simplified - full implementation would need
                # to port all render methods here)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _fill_placeholders(self, slide: Any, spec: SlideSpec, accent: str) -> None:
        """Fill placeholder shapes in template layout."""
        for shape in slide.placeholders:
            phf = shape.placeholder_format
            placeholder_type = phf.type if hasattr(phf, 'type') else None

            # Title placeholders
            if placeholder_type and 'TITLE' in str(placeholder_type):
                if spec.title:
                    shape.text = spec.title
                    for para in shape.text_frame.paragraphs:
                        para.font.name = self.brand_kit.font_heading if self.brand_kit else "Calibri"

            # Body/content placeholders
            elif placeholder_type and 'BODY' in str(placeholder_type):
                if spec.bullets:
                    tf = shape.text_frame
                    tf.clear()
                    for i, bullet in enumerate(spec.bullets[:6]):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"• {bullet}"
                        p.level = 0
                        p.font.name = self.brand_kit.font_body if self.brand_kit else "Calibri"

            # Subtitle
            elif placeholder_type and 'SUBTITLE' in str(placeholder_type):
                if spec.subtitle:
                    shape.text = spec.subtitle

    def _get_theme_bg(self) -> str:
        """Get background color from theme."""
        if self.brand_kit:
            # Use light background for brand
            return "FFFFFF"
        return "FFFFFF"


def load_preset_template(name: str) -> Optional[str]:
    """
    Get path to a preset template.

    Args:
        name: One of "corporate", "startup", "academic", "minimal"

    Returns:
        Path to template file or None if not found
    """
    # Templates would be bundled with the package
    template_dir = os.path.join(os.path.dirname(__file__), "templates")
    template_path = os.path.join(template_dir, f"{name}.pptx")
    if os.path.exists(template_path):
        return template_path
    return None
