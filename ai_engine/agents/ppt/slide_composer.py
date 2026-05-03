"""
SlideComposer — deterministic DeckSpec → .pptx assembler.

Uses python-pptx. Builds slides from scratch (no template file required) so
the package works in any environment. Themes map to (background, accent,
title, body) color quartets and a font family.

Layout decisions are deterministic per (kind, available content) — no LLM in
this layer. Charts and images are optional: if absent, a placeholder card is
drawn so the slide still looks intentional.

Public API:
    SlideComposer().compose(deck) -> bytes  (pptx file content)
"""
from __future__ import annotations

import io
import logging
from typing import Dict, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ai_engine.agents.ppt.schemas import DeckSpec, SlideKind, SlideSpec

logger = logging.getLogger(__name__)


# ────────────────────────────────────────────────────────────────────────
#  Themes
# ────────────────────────────────────────────────────────────────────────

#: (background, accent, title, body) — RGB hex strings.
_Theme = Tuple[str, str, str, str]

THEMES: Dict[str, _Theme] = {
    "modern":     ("FFFFFF", "2563EB", "0F172A", "334155"),
    "midnight":   ("0F172A", "38BDF8", "F8FAFC", "CBD5E1"),
    "warm":       ("FFFBEB", "EA580C", "1F2937", "374151"),
    "minimal":    ("FAFAFA", "111827", "111827", "4B5563"),
    "vibrant":    ("FFFFFF", "DB2777", "111827", "374151"),
    "corporate":  ("FFFFFF", "0F4C81", "0B2A4A", "334E68"),
}

DEFAULT_FONT = "Calibri"

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


def _hex_to_rgb(value: str) -> RGBColor:
    h = value.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Bad hex color: {value!r}")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _resolve_theme(deck: DeckSpec) -> _Theme:
    base = THEMES.get(deck.theme.lower(), THEMES["modern"])
    if deck.accent_color:
        try:
            _hex_to_rgb(deck.accent_color)  # validate
            return (base[0], deck.accent_color.lstrip("#"), base[2], base[3])
        except ValueError:
            logger.debug("ppt_bad_accent: %s — using theme default", deck.accent_color)
    return base


# ────────────────────────────────────────────────────────────────────────
#  SlideComposer
# ────────────────────────────────────────────────────────────────────────

class SlideComposer:
    """Render a DeckSpec to .pptx bytes."""

    def __init__(self, *, chart_renderer: Optional[object] = None,
                 image_resolver: Optional[object] = None) -> None:
        # P2/P3 dependency-injection hooks. Optional — composer renders text-only
        # decks without them, drawing accent placeholders for charts/images.
        self.chart_renderer = chart_renderer
        self.image_resolver = image_resolver

    async def compose(self, deck: DeckSpec) -> bytes:
        bg, accent, title_c, body_c = _resolve_theme(deck)
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]  # truly blank

        for slide_spec in deck.slides:
            slide = prs.slides.add_slide(blank)
            self._paint_background(slide, bg)
            await self._render_slide(slide, slide_spec, accent=accent, title_c=title_c, body_c=body_c)
            self._add_speaker_notes(slide, slide_spec.notes)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ────────────────────────────────────────────────────────────────────
    #  Per-slide dispatch
    # ────────────────────────────────────────────────────────────────────

    async def _render_slide(
        self, slide, spec: SlideSpec, *, accent: str, title_c: str, body_c: str,
    ) -> None:
        kind = spec.kind
        if kind == SlideKind.title:
            self._render_title(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind == SlideKind.section:
            self._render_section(slide, spec, accent=accent, title_c=title_c)
        elif kind == SlideKind.two_column:
            self._render_two_column(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind == SlideKind.quote:
            self._render_quote(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind == SlideKind.chart:
            await self._render_chart_slide(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind in (SlideKind.image, SlideKind.image_text):
            await self._render_image_slide(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind == SlideKind.table:
            self._render_table(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        elif kind == SlideKind.closing:
            self._render_closing(slide, spec, accent=accent, title_c=title_c, body_c=body_c)
        else:  # content (default)
            self._render_content(slide, spec, accent=accent, title_c=title_c, body_c=body_c)

    # ────────────────────────────────────────────────────────────────────
    #  Background + decoration
    # ────────────────────────────────────────────────────────────────────

    def _paint_background(self, slide, bg_hex: str) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(bg_hex)

    def _add_accent_bar(self, slide, accent: str, *, top: Emu, height: Emu = Inches(0.08)) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), top, Inches(1.5), height,
        )
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = _hex_to_rgb(accent)

    def _add_speaker_notes(self, slide, notes: Optional[str]) -> None:
        if not notes:
            return
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes

    # ────────────────────────────────────────────────────────────────────
    #  Helpers: text boxes
    # ────────────────────────────────────────────────────────────────────

    def _add_text(
        self, slide, *, left, top, width, height, text: str,
        size: int = 24, bold: bool = False, color: str = "111827",
        align=PP_ALIGN.LEFT,
    ):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = DEFAULT_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _hex_to_rgb(color)
        return tb

    def _add_bullets(
        self, slide, *, left, top, width, height, bullets, color: str = "334155", size: int = 18,
    ):
        if not bullets:
            return None
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.name = DEFAULT_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = _hex_to_rgb(color)
        return tb

    # ────────────────────────────────────────────────────────────────────
    #  Layouts
    # ────────────────────────────────────────────────────────────────────

    def _render_title(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(2.6), height=Inches(0.12))
        self._add_text(slide, left=Inches(0.7), top=Inches(2.85),
                       width=Inches(11.9), height=Inches(1.6),
                       text=spec.title or "", size=44, bold=True, color=title_c)
        if spec.subtitle:
            self._add_text(slide, left=Inches(0.7), top=Inches(4.4),
                           width=Inches(11.9), height=Inches(0.8),
                           text=spec.subtitle, size=22, color=body_c)

    def _render_section(self, slide, spec: SlideSpec, *, accent, title_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(3.4), height=Inches(0.18))
        self._add_text(slide, left=Inches(0.7), top=Inches(3.6),
                       width=Inches(11.9), height=Inches(1.4),
                       text=spec.title or "", size=40, bold=True, color=title_c)

    def _render_content(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(0.55))
        self._add_text(slide, left=Inches(0.7), top=Inches(0.75),
                       width=Inches(11.9), height=Inches(0.9),
                       text=spec.title or "", size=30, bold=True, color=title_c)
        if spec.subtitle:
            self._add_text(slide, left=Inches(0.7), top=Inches(1.55),
                           width=Inches(11.9), height=Inches(0.5),
                           text=spec.subtitle, size=16, color=body_c)
            bullet_top = Inches(2.2)
        else:
            bullet_top = Inches(1.8)
        self._add_bullets(slide, left=Inches(0.9), top=bullet_top,
                          width=Inches(11.5), height=Inches(5.0),
                          bullets=spec.bullets, color=body_c, size=20)

    def _render_two_column(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(0.55))
        self._add_text(slide, left=Inches(0.7), top=Inches(0.75),
                       width=Inches(11.9), height=Inches(0.9),
                       text=spec.title or "", size=30, bold=True, color=title_c)
        self._add_bullets(slide, left=Inches(0.9), top=Inches(1.9),
                          width=Inches(5.6), height=Inches(5.0),
                          bullets=spec.bullets, color=body_c, size=18)
        self._add_bullets(slide, left=Inches(6.9), top=Inches(1.9),
                          width=Inches(5.6), height=Inches(5.0),
                          bullets=spec.bullets_right or [], color=body_c, size=18)

    def _render_quote(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(2.4), height=Inches(0.14))
        body = spec.body or spec.title or ""
        self._add_text(slide, left=Inches(1.2), top=Inches(2.7),
                       width=Inches(10.9), height=Inches(2.6),
                       text=f"“{body}”", size=32, bold=False, color=title_c)
        if spec.attribution:
            self._add_text(slide, left=Inches(1.2), top=Inches(5.4),
                           width=Inches(10.9), height=Inches(0.6),
                           text=f"— {spec.attribution}", size=18, color=body_c)

    def _render_table(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        """Render a table slide."""
        self._add_accent_bar(slide, accent, top=Inches(0.55))
        self._add_text(slide, left=Inches(0.7), top=Inches(0.75),
                       width=Inches(11.9), height=Inches(0.9),
                       text=spec.title or "", size=28, bold=True, color=title_c)

        table_spec = spec.table
        if table_spec is None:
            self._draw_placeholder(slide, accent, label="Table",
                                   left=Inches(1.0), top=Inches(1.9),
                                   width=Inches(11.3), height=Inches(4.6))
            return

        # Table dimensions
        rows = len(table_spec.rows) + 1  # +1 for header
        cols = len(table_spec.headers) if table_spec.headers else (len(table_spec.rows[0]) if table_spec.rows else 2)
        if cols == 0:
            cols = 2

        # Position table below title
        table_left = Inches(0.7)
        table_top = Inches(1.8)
        table_width = Inches(11.9)
        table_height = Inches(5.0)

        try:
            table = slide.shapes.add_table(
                rows, cols,
                table_left, table_top, table_width, table_height,
            ).table

            # Style based on spec
            style_map = {
                "light": 1,  # Table Style Light 1
                "medium": 2,  # Table Style Medium 2
                "dark": 3,  # Table Style Dark 1
            }
            table_style = style_map.get(table_spec.style, 2)

            # Fill headers
            if table_spec.headers:
                for col_idx, header in enumerate(table_spec.headers[:cols]):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    # Header styling
                    for para in cell.text_frame.paragraphs:
                        para.font.bold = True
                        para.font.size = Pt(12)
                        para.font.name = DEFAULT_FONT
                        para.font.color.rgb = _hex_to_rgb(title_c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(accent)
                    cell.font.color.rgb = _hex_to_rgb("FFFFFF")

            # Fill data rows
            for row_idx, row_data in enumerate(table_spec.rows[:rows-1]):
                for col_idx, cell_text in enumerate(row_data[:cols]):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_text)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(11)
                        para.font.name = DEFAULT_FONT
                        para.font.color.rgb = _hex_to_rgb(body_c)

        except Exception as exc:  # noqa: BLE001
            logger.warning("table_render_failed: %s", str(exc)[:200])
            self._draw_placeholder(slide, accent, label="Table",
                                   left=Inches(1.0), top=Inches(1.9),
                                   width=Inches(11.3), height=Inches(4.6))

    async def _render_chart_slide(
        self, slide, spec: SlideSpec, *, accent, title_c, body_c,
    ) -> None:
        self._add_accent_bar(slide, accent, top=Inches(0.55))
        self._add_text(slide, left=Inches(0.7), top=Inches(0.75),
                       width=Inches(11.9), height=Inches(0.9),
                       text=spec.title or "", size=28, bold=True, color=title_c)

        # PRIORITY 1: Try native editable chart (users can right-click → Edit Data)
        rendered = False
        if spec.chart is not None:
            try:
                from ai_engine.agents.ppt.native_chart_renderer import NativeChartRenderer
                native_renderer = NativeChartRenderer()
                if native_renderer.can_render(spec.chart):
                    chart = native_renderer.render_to_slide(
                        slide, spec.chart,
                        accent_hex=accent,
                        left=Inches(1.0), top=Inches(1.9),
                        width=Inches(11.3), height=Inches(4.6),
                    )
                    if chart is not None:
                        rendered = True
                        logger.debug("ppt_native_chart_rendered: %s", spec.chart.kind)
            except Exception as exc:  # noqa: BLE001
                logger.debug("ppt_native_chart_failed: %s — falling back", str(exc)[:150])

        # PRIORITY 2: Fall back to matplotlib PNG (static image)
        if not rendered and self.chart_renderer is not None and spec.chart is not None:
            try:
                png = await self.chart_renderer.render(spec.chart, accent_hex=accent)  # type: ignore[attr-defined]
                if png:
                    slide.shapes.add_picture(
                        io.BytesIO(png),
                        Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.6),
                    )
                    rendered = True
            except Exception as exc:  # noqa: BLE001
                logger.warning("ppt_chart_render_failed: %s", str(exc)[:200])

        # PRIORITY 3: Placeholder if all renderers fail
        if not rendered:
            self._draw_placeholder(slide, accent, label="Chart",
                                   left=Inches(1.0), top=Inches(1.9),
                                   width=Inches(11.3), height=Inches(4.6))
        if spec.caption:
            self._add_text(slide, left=Inches(0.7), top=Inches(6.6),
                           width=Inches(11.9), height=Inches(0.6),
                           text=spec.caption, size=14, color=body_c, align=PP_ALIGN.CENTER)

    async def _render_image_slide(
        self, slide, spec: SlideSpec, *, accent, title_c, body_c,
    ) -> None:
        self._add_accent_bar(slide, accent, top=Inches(0.55))
        self._add_text(slide, left=Inches(0.7), top=Inches(0.75),
                       width=Inches(11.9), height=Inches(0.9),
                       text=spec.title or "", size=28, bold=True, color=title_c)
        if spec.kind == SlideKind.image_text:
            img_box = (Inches(0.7), Inches(1.9), Inches(6.0), Inches(4.6))
            text_box = (Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.6))
        else:
            img_box = (Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.6))
            text_box = None
        await self._place_image(slide, spec, img_box=img_box, accent=accent)
        if text_box and spec.bullets:
            l, t, w, h = text_box
            self._add_bullets(slide, left=l, top=t, width=w, height=h,
                              bullets=spec.bullets, color=body_c, size=18)
        if spec.caption:
            self._add_text(slide, left=Inches(0.7), top=Inches(6.6),
                           width=Inches(11.9), height=Inches(0.6),
                           text=spec.caption, size=12, color=body_c, align=PP_ALIGN.CENTER)

    async def _place_image(self, slide, spec: SlideSpec, *, img_box, accent: str) -> None:
        l, t, w, h = img_box
        png_bytes: Optional[bytes] = None
        if spec.image is not None:
            if spec.image.bytes_b64:
                import base64
                try:
                    png_bytes = base64.b64decode(spec.image.bytes_b64)
                except Exception as exc:  # noqa: BLE001
                    logger.debug("ppt_image_b64_decode_failed: %s", str(exc)[:200])
            if png_bytes is None and self.image_resolver is not None:
                try:
                    png_bytes = await self.image_resolver.resolve(spec.image)  # type: ignore[attr-defined]
                except Exception as exc:  # noqa: BLE001
                    logger.warning("ppt_image_resolve_failed: %s", str(exc)[:200])
        if png_bytes:
            slide.shapes.add_picture(io.BytesIO(png_bytes), l, t, w, h)
        else:
            self._draw_placeholder(slide, accent, label="Image",
                                   left=l, top=t, width=w, height=h)

    def _render_closing(self, slide, spec: SlideSpec, *, accent, title_c, body_c) -> None:
        self._add_accent_bar(slide, accent, top=Inches(2.6), height=Inches(0.12))
        self._add_text(slide, left=Inches(0.7), top=Inches(2.85),
                       width=Inches(11.9), height=Inches(1.6),
                       text=spec.title or "Thank you", size=44, bold=True,
                       color=title_c, align=PP_ALIGN.CENTER)
        if spec.body:
            self._add_text(slide, left=Inches(0.7), top=Inches(4.5),
                           width=Inches(11.9), height=Inches(0.8),
                           text=spec.body, size=22, color=body_c, align=PP_ALIGN.CENTER)

    def _draw_placeholder(self, slide, accent: str, *, label: str,
                          left, top, width, height) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
        shape.line.color.rgb = _hex_to_rgb(accent)
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.text = f"[{label}]"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = _hex_to_rgb(accent)
                run.font.name = DEFAULT_FONT
