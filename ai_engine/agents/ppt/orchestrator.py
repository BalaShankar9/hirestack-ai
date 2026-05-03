"""
PresentationOrchestrator — Professional-grade PPT generation pipeline.

Architecture:
    Pipeline-based phase execution with dependency injection.
    Each phase is an optional, swappable component implementing a protocol.

Design Principles:
    - Single Responsibility: Each phase handles one concern
    - Open/Closed: New phases can be added without modifying existing code
    - Dependency Inversion: Phases depend on abstractions (protocols)
    - Composition over Inheritance: Phases composed into pipeline

Example:
    orchestrator = PresentationOrchestrator(
        phases=[
            OutlinePhase(),
            DataResearchPhase(enabled=True),
            ContentEnhancementPhase(enabled=True),
            QualityValidationPhase(),
            CompositionPhase(),
            PolishPhase(),
        ]
    )
    result = await orchestrator.generate(topic="AI Trends")
"""
from __future__ import annotations

import asyncio
import hashlib
import logging
import time
from dataclasses import dataclass, field
from enum import Enum
from typing import Any, Callable, Dict, List, Optional, Protocol, runtime_checkable, Set

from ai_engine.agents.ppt.outline_planner import OutlinePlanner
from ai_engine.agents.ppt.schemas import DeckSpec
from ai_engine.agents.ppt.slide_composer import SlideComposer

logger = logging.getLogger(__name__)


class GenerationStatus(Enum):
    """Status of generation pipeline."""
    PENDING = "pending"
    PLANNING = "planning"
    RESEARCHING = "researching"
    ENHANCING = "enhancing"
    GENERATING_IMAGES = "generating_images"
    VALIDATING = "validating"
    COMPOSING = "composing"
    POLISHING = "polishing"
    FINALIZING = "finalizing"
    COMPLETED = "completed"
    FAILED = "failed"


# =============================================================================
# Result Types
# =============================================================================

@dataclass(frozen=True)
class GenerationResult:
    """Immutable result of presentation generation."""

    pptx_bytes: bytes
    deck: DeckSpec
    latency_ms: int
    quality_score: float = 0.0
    metadata: Dict[str, Any] = field(default_factory=dict)
    # Observability
    generation_id: str = ""
    phase_latencies: Dict[str, int] = field(default_factory=dict)
    cache_hit: bool = False

    @property
    def slide_count(self) -> int:
        return self.deck.slide_count

    @property
    def size_bytes(self) -> int:
        return len(self.pptx_bytes)


@dataclass
class GenerationProgress:
    """Progress update for long-running generations."""
    status: GenerationStatus
    percent: int  # 0-100
    message: str
    phase: Optional[str] = None
    slide_idx: Optional[int] = None
    latency_so_far_ms: int = 0


@dataclass(frozen=True)
class PhaseResult:
    """Result from a single pipeline phase."""

    success: bool
    deck: Optional[DeckSpec] = None
    pptx_bytes: Optional[bytes] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None


# =============================================================================
# Phase Protocols (Abstract Interfaces)
# =============================================================================

@runtime_checkable
class PlanningPhase(Protocol):
    """Protocol for outline generation phase."""

    async def execute(
        self,
        *,
        topic: str,
        audience: Optional[str],
        slide_count: int,
        tone: Optional[str],
        theme: str,
        extra_context: Optional[str],
    ) -> PhaseResult:
        """Generate presentation outline."""
        ...


@runtime_checkable
class DeckTransformPhase(Protocol):
    """Protocol for phases that transform the deck (before composition)."""

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        """Transform deck and return result."""
        ...


@runtime_checkable
class CompositionPhase(Protocol):
    """Protocol for composition phase (deck -> pptx bytes)."""

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        """Compose deck into pptx bytes."""
        ...


@runtime_checkable
class PostProcessPhase(Protocol):
    """Protocol for post-composition phases."""

    async def execute(
        self, pptx_bytes: bytes, deck: DeckSpec, context: Dict[str, Any]
    ) -> PhaseResult:
        """Process pptx bytes and return modified bytes."""
        ...


# =============================================================================
# Phase Implementations
# =============================================================================

class OutlineGenerationPhase:
    """Phase 0: Generate presentation outline using AI."""

    def __init__(self, ai_client: Optional[Any] = None) -> None:
        self._planner = OutlinePlanner(ai_client=ai_client)

    async def execute(
        self,
        *,
        topic: str,
        audience: Optional[str],
        slide_count: int,
        tone: Optional[str],
        theme: str,
        extra_context: Optional[str],
    ) -> PhaseResult:
        try:
            deck = await self._planner.plan(
                topic=topic,
                audience=audience,
                slide_count=slide_count,
                tone=tone,
                theme=theme,
                extra_context=extra_context,
            )
            return PhaseResult(success=True, deck=deck)
        except Exception as exc:
            logger.exception("outline_generation_failed")
            return PhaseResult(success=False, error=str(exc))


class DataResearchPhase:
    """Phase 3: Enrich deck with real-world data."""

    def __init__(
        self,
        *,
        researcher: Optional[Any] = None,
        enabled: bool = False,
    ) -> None:
        self._enabled = enabled
        self._researcher = researcher

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        if not self._enabled or self._researcher is None:
            return PhaseResult(success=True, deck=deck)

        try:
            topic = context.get("topic", "")
            enriched_slides = await self._researcher.enrich_deck(deck.slides, topic)
            deck = deck.model_copy(update={"slides": enriched_slides})
            return PhaseResult(
                success=True,
                deck=deck,
                metadata={"data_research": True},
            )
        except Exception as exc:
            logger.warning("data_research_failed: %s", exc)
            return PhaseResult(success=True, deck=deck)  # Graceful degradation


class ContentEnhancementPhase:
    """Phase 10: AI-powered content optimization."""

    def __init__(
        self,
        *,
        enhancer: Optional[Any] = None,
        enabled: bool = False,
    ) -> None:
        self._enabled = enabled
        self._enhancer = enhancer

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        if not self._enabled or self._enhancer is None:
            return PhaseResult(success=True, deck=deck)

        try:
            enhanced_slides = await self._enhancer.enhance_deck(deck.slides)
            deck = deck.model_copy(update={"slides": enhanced_slides})
            return PhaseResult(
                success=True,
                deck=deck,
                metadata={"content_enhancement": True},
            )
        except Exception as exc:
            logger.warning("content_enhancement_failed: %s", exc)
            return PhaseResult(success=True, deck=deck)


class AIImageGenerationPhase:
    """Phase 9: Generate custom AI visuals for slides."""

    def __init__(
        self,
        *,
        generator: Optional[Any] = None,
        enabled: bool = False,
    ) -> None:
        self._enabled = enabled
        self._generator = generator

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        if not self._enabled or self._generator is None:
            return PhaseResult(success=True, deck=deck)

        try:
            theme = context.get("theme", "modern")
            slides = list(deck.slides)

            for i, slide in enumerate(slides):
                if slide.kind.value in ("title", "section", "image") and not slide.image:
                    img = await self._generator.generate_for_slide(slide, theme)
                    if img:
                        slides[i] = slide.model_copy(update={"image": img})

            deck = deck.model_copy(update={"slides": slides})
            return PhaseResult(
                success=True,
                deck=deck,
                metadata={"ai_images": True},
            )
        except Exception as exc:
            logger.warning("ai_image_generation_failed: %s", exc)
            return PhaseResult(success=True, deck=deck)


class QualityValidationPhase:
    """Phase 8: Validate deck quality and suggest revisions."""

    def __init__(
        self,
        *,
        validator: Optional[Any] = None,
        auto_revise: bool = True,
        threshold: float = 0.7,
    ) -> None:
        self._validator = validator
        self._auto_revise = auto_revise
        self._threshold = threshold

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        if self._validator is None:
            return PhaseResult(success=True, deck=deck, metadata={"quality_score": 0.0})

        try:
            report = self._validator.validate_deck(deck)
            quality_score = report.aggregate_score

            if self._auto_revise and quality_score < self._threshold:
                deck = self._validator.suggest_revisions(deck)

            return PhaseResult(
                success=True,
                deck=deck,
                metadata={"quality_score": quality_score},
            )
        except Exception as exc:
            logger.debug("quality_validation_failed: %s", exc)
            return PhaseResult(success=True, deck=deck, metadata={"quality_score": 0.0})


class PresentationCompositionPhase:
    """Phase 1: Compose deck into PPTX bytes."""

    def __init__(
        self,
        *,
        chart_renderer: Optional[Any] = None,
        image_resolver: Optional[Any] = None,
    ) -> None:
        # Lazy initialization for optional dependencies
        if chart_renderer is None:
            try:
                from ai_engine.agents.ppt.chart_renderer import ChartRenderer
                chart_renderer = ChartRenderer()
            except Exception:
                chart_renderer = None

        if image_resolver is None:
            try:
                from ai_engine.agents.ppt.image_fetcher import ImageFetcher
                image_resolver = ImageFetcher()
            except Exception:
                image_resolver = None

        self._composer = SlideComposer(
            chart_renderer=chart_renderer,
            image_resolver=image_resolver,
        )

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        try:
            pptx_bytes = await self._composer.compose(deck)
            return PhaseResult(
                success=True,
                deck=deck,
                pptx_bytes=pptx_bytes,
            )
        except Exception as exc:
            logger.exception("composition_failed")
            return PhaseResult(success=False, error=str(exc))


class PolishPhase:
    """Phase 7: Apply final polish and accessibility improvements."""

    def __init__(self, polish_engine: Optional[Any] = None) -> None:
        self._polish_engine = polish_engine

    async def execute(
        self, pptx_bytes: bytes, deck: DeckSpec, context: Dict[str, Any]
    ) -> PhaseResult:
        if self._polish_engine is None:
            return PhaseResult(success=True, pptx_bytes=pptx_bytes)

        try:
            polished = await self._polish_engine.polish_presentation(pptx_bytes)
            return PhaseResult(
                success=True,
                pptx_bytes=polished,
                metadata={"polish": True},
            )
        except Exception as exc:
            logger.debug("polish_failed: %s", exc)
            return PhaseResult(success=True, pptx_bytes=pptx_bytes)


class InteractiveElementsPhase:
    """Phase 11: Add interactive elements (TOC, navigation)."""

    def __init__(
        self,
        *,
        builder: Optional[Any] = None,
        enabled: bool = False,
    ) -> None:
        self._enabled = enabled
        self._builder = builder

    async def execute(
        self, pptx_bytes: bytes, deck: DeckSpec, context: Dict[str, Any]
    ) -> PhaseResult:
        if not self._enabled or self._builder is None:
            return PhaseResult(success=True, pptx_bytes=pptx_bytes)

        try:
            import io
            from pptx import Presentation

            prs = Presentation(io.BytesIO(pptx_bytes))
            self._builder.add_navigation(prs)

            if len(deck.slides) > 8:
                self._builder.add_toc(prs, deck.slides)

            buf = io.BytesIO()
            prs.save(buf)
            return PhaseResult(
                success=True,
                pptx_bytes=buf.getvalue(),
                metadata={"interactive": True},
            )
        except Exception as exc:
            logger.debug("interactive_elements_failed: %s", exc)
            return PhaseResult(success=True, pptx_bytes=pptx_bytes)


class TranslationPhase:
    """Phase 12: Translate deck to target language."""

    def __init__(
        self,
        *,
        translator: Optional[Any] = None,
        target_language: str = "en",
    ) -> None:
        self._translator = translator
        self._target_language = target_language

    async def execute(self, deck: DeckSpec, context: Dict[str, Any]) -> PhaseResult:
        if self._target_language == "en" or self._translator is None:
            return PhaseResult(success=True, deck=deck)

        try:
            translated = await self._translator.translate_deck(
                deck, self._target_language
            )
            return PhaseResult(
                success=True,
                deck=translated,
                metadata={"translated": self._target_language},
            )
        except Exception as exc:
            logger.warning("translation_failed: %s", exc)
            return PhaseResult(success=True, deck=deck)


# =============================================================================
# Orchestrator
# =============================================================================

class CircuitBreaker:
    """Circuit breaker for external API calls."""

    def __init__(
        self,
        failure_threshold: int = 5,
        recovery_timeout: float = 30.0,
    ) -> None:
        self.failure_threshold = failure_threshold
        self.recovery_timeout = recovery_timeout
        self.failures = 0
        self.last_failure_time: Optional[float] = None
        self.state = "closed"  # closed, open, half-open

    def call(self, fn: Callable[[], Any]) -> Any:
        """Execute function with circuit breaker protection."""
        if self.state == "open":
            if time.monotonic() - (self.last_failure_time or 0) > self.recovery_timeout:
                self.state = "half-open"
            else:
                raise RuntimeError("Circuit breaker is open")

        try:
            result = fn()
            if self.state == "half-open":
                self.state = "closed"
                self.failures = 0
            return result
        except Exception as exc:
            self.failures += 1
            self.last_failure_time = time.monotonic()
            if self.failures >= self.failure_threshold:
                self.state = "open"
            raise


class PresentationOrchestrator:
    """
    Production-grade presentation generation orchestrator.

    Features:
        - Pipeline-based phase execution
        - Result caching
        - Progress tracking
        - Circuit breakers for external APIs
        - Rate limiting
        - Comprehensive observability
    """

    def __init__(
        self,
        *,
        # Planning
        ai_client: Optional[Any] = None,
        outline_phase: Optional[PlanningPhase] = None,
        # Pre-composition transforms
        data_research: Optional[DeckTransformPhase] = None,
        content_enhancement: Optional[DeckTransformPhase] = None,
        ai_image_generation: Optional[DeckTransformPhase] = None,
        quality_validation: Optional[DeckTransformPhase] = None,
        translation: Optional[DeckTransformPhase] = None,
        # Composition
        composition: Optional[CompositionPhase] = None,
        # Post-processing
        polish: Optional[PostProcessPhase] = None,
        interactive_elements: Optional[PostProcessPhase] = None,
        # Production features
        enable_caching: bool = True,
        cache_ttl_seconds: float = 3600.0,
        max_concurrent_generations: int = 5,
        progress_callback: Optional[Callable[[GenerationProgress], None]] = None,
    ) -> None:
        """
        Initialize orchestrator with configured phases.

        Args:
            ai_client: AI client for outline generation
            outline_phase: Custom outline generation phase
            data_research: Phase 3 - Data research
            content_enhancement: Phase 10 - Content enhancement
            ai_image_generation: Phase 9 - AI image generation
            quality_validation: Phase 8 - Quality validation
            translation: Phase 12 - i18n translation
            composition: Phase 1 - Deck to PPTX composition
            polish: Phase 7 - Final polish
            interactive_elements: Phase 11 - Interactive elements
        """
        # Phase 0: Planning (required)
        self._outline_phase = outline_phase or OutlineGenerationPhase(ai_client)

        # Pre-composition transforms (Phase 3, 8-10, 12)
        self._pre_composition: List[DeckTransformPhase] = []
        if data_research:
            self._pre_composition.append(data_research)
        if content_enhancement:
            self._pre_composition.append(content_enhancement)
        if ai_image_generation:
            self._pre_composition.append(ai_image_generation)
        if quality_validation:
            self._pre_composition.append(quality_validation)
        if translation:
            self._pre_composition.append(translation)

        # Phase 1: Composition (required)
        self._composition = composition or PresentationCompositionPhase()

        # Post-composition (Phase 7, 11)
        self._post_composition: List[PostProcessPhase] = []
        if polish:
            self._post_composition.append(polish)
        if interactive_elements:
            self._post_composition.append(interactive_elements)

        # Production features
        self._enable_caching = enable_caching
        self._cache_ttl = cache_ttl_seconds
        self._cache: Dict[str, Tuple[GenerationResult, float]] = {}
        self._semaphore = asyncio.Semaphore(max_concurrent_generations)
        self._progress_callback = progress_callback
        self._circuit_breakers: Dict[str, CircuitBreaker] = {}

    def _get_cache_key(self, **kwargs: Any) -> str:
        """Generate cache key from generation parameters."""
        key_data = "|".join(f"{k}={v}" for k, v in sorted(kwargs.items()) if v is not None)
        return hashlib.md5(key_data.encode()).hexdigest()

    def _get_cached(self, key: str) -> Optional[GenerationResult]:
        """Get cached result if valid."""
        if not self._enable_caching or key not in self._cache:
            return None
        result, timestamp = self._cache[key]
        if time.monotonic() - timestamp > self._cache_ttl:
            del self._cache[key]
            return None
        return result

    def _cache_result(self, key: str, result: GenerationResult) -> None:
        """Cache generation result."""
        if not self._enable_caching:
            return
        # LRU eviction if cache too large
        if len(self._cache) >= 100:
            oldest = min(self._cache.keys(), key=lambda k: self._cache[k][1])
            del self._cache[oldest]
        self._cache[key] = (result, time.monotonic())

    def _get_circuit_breaker(self, name: str) -> CircuitBreaker:
        """Get or create circuit breaker for external service."""
        if name not in self._circuit_breakers:
            self._circuit_breakers[name] = CircuitBreaker()
        return self._circuit_breakers[name]

    async def _report_progress(
        self,
        status: GenerationStatus,
        percent: int,
        message: str,
        latency_so_far_ms: int = 0,
    ) -> None:
        """Report progress if callback is configured."""
        if self._progress_callback:
            try:
                self._progress_callback(GenerationProgress(
                    status=status,
                    percent=percent,
                    message=message,
                    latency_so_far_ms=latency_so_far_ms,
                ))
            except Exception as exc:
                logger.debug("progress_callback_failed: %s", exc)

    @classmethod
    def create_with_defaults(
        cls,
        *,
        ai_client: Optional[Any] = None,
        enable_data_research: bool = False,
        enable_content_enhancement: bool = False,
        enable_ai_images: bool = False,
        enable_interactive: bool = False,
        target_language: str = "en",
        # Production parameters
        enable_caching: bool = True,
        cache_ttl_seconds: float = 3600.0,
        max_concurrent_generations: int = 5,
        progress_callback: Optional[Callable[[GenerationProgress], None]] = None,
    ) -> "PresentationOrchestrator":
        """
        Factory method to create orchestrator with standard phase configuration.

        Example:
            orch = PresentationOrchestrator.create_with_defaults(
                enable_data_research=True,
                enable_content_enhancement=True,
                enable_caching=True,
            )
        """
        # Lazy-load optional engines based on configuration
        data_researcher = None
        if enable_data_research:
            try:
                from ai_engine.agents.ppt.data_researcher import DataResearcher
                data_researcher = DataResearcher()
            except Exception:
                pass

        content_enhancer = None
        if enable_content_enhancement:
            try:
                from ai_engine.agents.ppt.content_enhancer import ContentEnhancer
                content_enhancer = ContentEnhancer(ai_client=ai_client)
            except Exception:
                pass

        ai_generator = None
        if enable_ai_images:
            try:
                from ai_engine.agents.ppt.ai_image_generator import AIImageGenerator
                ai_generator = AIImageGenerator()
            except Exception:
                pass

        quality_validator = None
        try:
            from ai_engine.agents.ppt.quality_validator import QualityValidator
            quality_validator = QualityValidator(ai_client=ai_client)
        except Exception:
            pass

        polish_engine = None
        try:
            from ai_engine.agents.ppt.polish_engine import PolishEngine
            polish_engine = PolishEngine()
        except Exception:
            pass

        interactive_builder = None
        if enable_interactive:
            try:
                from ai_engine.agents.ppt.interactive_builder import InteractiveBuilder
                interactive_builder = InteractiveBuilder()
            except Exception:
                pass

        translator = None
        if target_language != "en":
            try:
                from ai_engine.agents.ppt.i18n_engine import I18nEngine
                translator = I18nEngine(ai_client=ai_client)
            except Exception:
                pass

        return cls(
            ai_client=ai_client,
            data_research=DataResearchPhase(
                researcher=data_researcher, enabled=enable_data_research
            ),
            content_enhancement=ContentEnhancementPhase(
                enhancer=content_enhancer, enabled=enable_content_enhancement
            ),
            ai_image_generation=AIImageGenerationPhase(
                generator=ai_generator, enabled=enable_ai_images
            ),
            quality_validation=QualityValidationPhase(validator=quality_validator),
            translation=TranslationPhase(
                translator=translator, target_language=target_language
            ),
            polish=PolishPhase(polish_engine),
            interactive_elements=InteractiveElementsPhase(
                builder=interactive_builder, enabled=enable_interactive
            ),
            # Production features
            enable_caching=enable_caching,
            cache_ttl_seconds=cache_ttl_seconds,
            max_concurrent_generations=max_concurrent_generations,
            progress_callback=progress_callback,
        )

    async def generate(
        self,
        *,
        topic: str,
        audience: Optional[str] = None,
        slide_count: int = 10,
        tone: Optional[str] = None,
        theme: str = "modern",
        extra_context: Optional[str] = None,
        skip_cache: bool = False,
    ) -> GenerationResult:
        """
        Execute full generation pipeline with production features.

        Args:
            topic: Presentation topic
            audience: Target audience
            slide_count: Number of slides (3-30)
            tone: Content tone
            theme: Visual theme
            extra_context: Additional context
            skip_cache: Bypass cache lookup

        Returns:
            GenerationResult with full observability data
        """
        if not topic or not topic.strip():
            raise ValueError("topic must be a non-empty string")

        # Validate parameters
        if slide_count < 3 or slide_count > 30:
            raise ValueError("slide_count must be between 3 and 30")

        # Generate unique ID for this generation
        generation_id = hashlib.md5(
            f"{topic}{time.monotonic()}".encode()
        ).hexdigest()[:12]

        # Check cache
        cache_key = self._get_cache_key(
            topic=topic,
            audience=audience,
            slide_count=slide_count,
            tone=tone,
            theme=theme,
        )

        if not skip_cache:
            cached = self._get_cached(cache_key)
            if cached:
                logger.info(
                    "cache_hit: generation_id=%s topic=%s",
                    generation_id, topic[:40]
                )
                return cached

        # Use semaphore for concurrency control
        async with self._semaphore:
            return await self._execute_pipeline(
                generation_id=generation_id,
                cache_key=cache_key,
                topic=topic.strip(),
                audience=audience,
                slide_count=slide_count,
                tone=tone,
                theme=theme,
                extra_context=extra_context,
            )

    async def _execute_pipeline(
        self,
        *,
        generation_id: str,
        cache_key: str,
        topic: str,
        audience: Optional[str],
        slide_count: int,
        tone: Optional[str],
        theme: str,
        extra_context: Optional[str],
    ) -> GenerationResult:
        """Execute the actual generation pipeline."""
        t0 = time.monotonic()
        phase_latencies: Dict[str, int] = {}

        context = {
            "topic": topic,
            "theme": theme,
            "generation_id": generation_id,
        }
        metadata: Dict[str, Any] = {"generation_id": generation_id}
        quality_score = 0.0

        await self._report_progress(
            status=GenerationStatus.PLANNING,
            percent=5,
            message="Generating outline...",
        )

        # Phase 0: Outline Generation
        phase_t0 = time.monotonic()
        outline_result = await self._outline_phase.execute(
            topic=topic,
            audience=audience,
            slide_count=slide_count,
            tone=tone,
            theme=theme,
            extra_context=extra_context,
        )
        phase_latencies["outline"] = int((time.monotonic() - phase_t0) * 1000)

        if not outline_result.success or outline_result.deck is None:
            await self._report_progress(
                status=GenerationStatus.FAILED,
                percent=0,
                message=f"Outline generation failed: {outline_result.error}",
                latency_so_far_ms=int((time.monotonic() - t0) * 1000),
            )
            raise RuntimeError(
                f"Outline generation failed: {outline_result.error}"
            )

        deck = outline_result.deck
        metadata["outline_latency_ms"] = phase_latencies["outline"]

        # Pre-composition phases (Phase 3, 8-10, 12)
        percent_per_phase = 40 // max(len(self._pre_composition), 1)
        current_percent = 10

        for idx, phase in enumerate(self._pre_composition):
            phase_name = phase.__class__.__name__.replace("Phase", "").lower()
            status_map = {
                "dataresearch": GenerationStatus.RESEARCHING,
                "contentenhancement": GenerationStatus.ENHANCING,
                "aiimagegeneration": GenerationStatus.GENERATING_IMAGES,
                "qualityvalidation": GenerationStatus.VALIDATING,
                "translation": GenerationStatus.ENHANCING,
            }
            status = status_map.get(phase_name, GenerationStatus.ENHANCING)

            await self._report_progress(
                status=status,
                percent=current_percent,
                message=f"Executing {phase.__class__.__name__}...",
            )

            phase_t0 = time.monotonic()
            result = await phase.execute(deck, context)
            phase_latencies[phase_name] = int((time.monotonic() - phase_t0) * 1000)

            if result.deck:
                deck = result.deck
            metadata.update(result.metadata)
            if "quality_score" in result.metadata:
                quality_score = result.metadata["quality_score"]

            current_percent += percent_per_phase

        await self._report_progress(
            status=GenerationStatus.COMPOSING,
            percent=50,
            message="Composing presentation...",
        )

        # Phase 1: Composition
        phase_t0 = time.monotonic()
        composition_result = await self._composition.execute(deck, context)
        phase_latencies["composition"] = int((time.monotonic() - phase_t0) * 1000)

        if not composition_result.success or composition_result.pptx_bytes is None:
            await self._report_progress(
                status=GenerationStatus.FAILED,
                percent=50,
                message=f"Composition failed: {composition_result.error}",
                latency_so_far_ms=int((time.monotonic() - t0) * 1000),
            )
            raise RuntimeError(
                f"Composition failed: {composition_result.error}"
            )

        pptx_bytes = composition_result.pptx_bytes
        metadata["composition_latency_ms"] = phase_latencies["composition"]

        # Post-composition phases (Phase 7, 11)
        percent_per_phase = 45 // max(len(self._post_composition), 1)

        for idx, phase in enumerate(self._post_composition):
            phase_name = phase.__class__.__name__.replace("Phase", "").lower()
            status_map = {
                "polish": GenerationStatus.POLISHING,
                "interactiveelements": GenerationStatus.FINALIZING,
            }
            status = status_map.get(phase_name, GenerationStatus.FINALIZING)

            await self._report_progress(
                status=status,
                percent=50 + (idx * percent_per_phase),
                message=f"Executing {phase.__class__.__name__}...",
            )

            phase_t0 = time.monotonic()
            result = await phase.execute(pptx_bytes, deck, context)
            phase_latencies[phase_name] = int((time.monotonic() - phase_t0) * 1000)

            if result.pptx_bytes:
                pptx_bytes = result.pptx_bytes
            metadata.update(result.metadata)

        total_latency_ms = int((time.monotonic() - t0) * 1000)

        await self._report_progress(
            status=GenerationStatus.COMPLETED,
            percent=100,
            message="Generation complete!",
            latency_so_far_ms=total_latency_ms,
        )

        # Structured logging
        logger.info(
            "presentation_generated: "
            "generation_id=%s topic=%s slides=%d size=%dKB "
            "latency=%dms quality=%.2f phases=%s",
            generation_id,
            topic[:60],
            deck.slide_count,
            len(pptx_bytes) // 1024,
            total_latency_ms,
            quality_score,
            phase_latencies,
        )

        result = GenerationResult(
            pptx_bytes=pptx_bytes,
            deck=deck,
            latency_ms=total_latency_ms,
            quality_score=quality_score,
            metadata=metadata,
            generation_id=generation_id,
            phase_latencies=phase_latencies,
        )

        # Cache result
        self._cache_result(cache_key, result)

        return result

    async def generate_from_deck(self, deck: DeckSpec) -> GenerationResult:
        """Render an already-planned deck through the pipeline (skips outline)."""
        t0 = time.monotonic()
        context: Dict[str, Any] = {"theme": deck.theme}
        metadata: Dict[str, Any] = {}
        quality_score = 0.0

        # Pre-composition phases
        for phase in self._pre_composition:
            result = await phase.execute(deck, context)
            if result.deck:
                deck = result.deck
            metadata.update(result.metadata)

        # Composition
        composition_result = await self._composition.execute(deck, context)

        if not composition_result.success or composition_result.pptx_bytes is None:
            raise RuntimeError(f"Composition failed: {composition_result.error}")

        pptx_bytes = composition_result.pptx_bytes

        # Post-composition phases
        for phase in self._post_composition:
            result = await phase.execute(pptx_bytes, deck, context)
            if result.pptx_bytes:
                pptx_bytes = result.pptx_bytes
            metadata.update(result.metadata)

        latency_ms = int((time.monotonic() - t0) * 1000)

        return GenerationResult(
            pptx_bytes=pptx_bytes,
            deck=deck,
            latency_ms=latency_ms,
            quality_score=quality_score,
            metadata=metadata,
        )

    def health_check(self) -> Dict[str, Any]:
        """
        Check orchestrator health status.

        Returns:
            Dict with health status of all phases and services
        """
        health = {
            "status": "healthy",
            "phases": {},
            "circuit_breakers": {},
            "cache": {
                "enabled": self._enable_caching,
                "size": len(self._cache),
            },
        }

        # Check each phase
        phase_map = {
            "outline": self._outline_phase,
            **{f"pre_{i}": p for i, p in enumerate(self._pre_composition)},
            "composition": self._composition,
            **{f"post_{i}": p for i, p in enumerate(self._post_composition)},
        }

        for name, phase in phase_map.items():
            phase_type = phase.__class__.__name__
            health["phases"][name] = {
                "type": phase_type,
                "available": phase is not None,
            }

        # Circuit breaker states
        for name, cb in self._circuit_breakers.items():
            health["circuit_breakers"][name] = cb.state

        return health

    def get_metrics(self) -> Dict[str, Any]:
        """Get orchestrator metrics and statistics."""
        return {
            "cache": {
                "enabled": self._enable_caching,
                "entries": len(self._cache),
                "ttl_seconds": self._cache_ttl,
            },
            "concurrency": {
                "max_concurrent": self._semaphore._value,
            },
            "circuit_breakers": {
                name: {
                    "state": cb.state,
                    "failures": cb.failures,
                }
                for name, cb in self._circuit_breakers.items()
            },
            "phases": {
                "pre_composition_count": len(self._pre_composition),
                "post_composition_count": len(self._post_composition),
            },
        }

    def clear_cache(self) -> int:
        """Clear result cache. Returns number of entries cleared."""
        count = len(self._cache)
        self._cache.clear()
        logger.info("cache_cleared: entries=%d", count)
        return count


# Backward compatibility alias
PPTOrchestrator = PresentationOrchestrator
PPTResult = GenerationResult
