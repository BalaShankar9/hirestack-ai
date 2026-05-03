"""
InteractiveBuilder — Hyperlinks, navigation, and engagement features.

Adds interactivity to presentations:
- Table of Contents with jump links
- External URL hyperlinks
- QR code generation for mobile access
- Slide-to-slide navigation buttons
- Action buttons ("Contact Us", "Learn More")
- Embedded video placeholders

Public API:
    InteractiveBuilder().add_toc(prs, slides) -> None
    InteractiveBuilder().add_hyperlink(shape, url) -> None
    InteractiveBuilder().generate_qr(data) -> bytes
    InteractiveBuilder().add_navigation(prs) -> None
    InteractiveBuilder().add_action_buttons(slide, actions) -> None
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

from ai_engine.agents.ppt.schemas import SlideKind, SlideSpec

logger = logging.getLogger(__name__)


@dataclass
class ActionButton:
    """Definition of an action button."""

    label: str
    action: str  # "url", "email", "slide", "file"
    target: str  # URL, email, slide number, or file path
    color: str = "#2563EB"
    position: Tuple[float, float] = (0.5, 6.5)  # (left, top) in inches
    size: Tuple[float, float] = (2.0, 0.5)  # (width, height) in inches


class InteractiveBuilder:
    """Add interactive elements to presentations."""

    # QR code error correction levels
    QR_ERROR_CORRECTION = {
        "L": 1,  # ~7% correction
        "M": 0,  # ~15% correction (default)
        "Q": 3,  # ~25% correction
        "H": 2,  # ~30% correction
    }

    def __init__(self) -> None:
        self._slide_anchors: Dict[str, int] = {}  # title -> slide index

    # ───────────────────────────────────────────────────────────────────────
    #  Table of Contents
    # ───────────────────────────────────────────────────────────────────────

    def add_toc(
        self,
        prs: Any,
        slides: List[SlideSpec],
        position: str = "after_title",
        include_slide_numbers: bool = True,
    ) -> int:
        """
        Add a clickable Table of Contents slide.

        Args:
            prs: Presentation object
            slides: Slide specifications
            position: "after_title", "before_title", "end", or slide index
            include_slide_numbers: Whether to show slide numbers

        Returns:
            Index of the added TOC slide
        """
        from pptx.util import Inches, Pt

        # Build TOC entries
        entries = []
        for idx, slide in enumerate(slides):
            if slide.kind in (SlideKind.section, SlideKind.title, SlideKind.content):
                if slide.title:
                    entries.append({
                        "title": slide.title,
                        "slide_idx": idx,
                        "kind": slide.kind,
                    })
                    # Store anchor
                    self._slide_anchors[slide.title.lower().replace(" ", "_")] = idx

        if len(entries) < 3:
            logger.debug("toc_skipped: insufficient entries (%d)", len(entries))
            return -1

        # Determine insert position
        insert_idx = 1 if position == "after_title" else 1
        if isinstance(position, int):
            insert_idx = position

        # Add TOC slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        toc_slide = prs.slides.add_slide(blank_layout)

        # Add title
        self._add_text_box(
            toc_slide,
            "Table of Contents",
            left=Inches(0.7),
            top=Inches(0.5),
            width=Inches(12),
            height=Inches(1),
            font_size=32,
            bold=True,
        )

        # Add entries
        y_pos = 1.5
        for entry in entries[:10]:  # Max 10 entries
            text = entry["title"]
            if include_slide_numbers:
                text = f"{entry['slide_idx'] + 1}. {text}"

            # Add clickable text (simulated with shape)
            shape = self._add_text_box(
                toc_slide,
                text,
                left=Inches(1.0),
                top=Inches(y_pos),
                width=Inches(11),
                height=Inches(0.6),
                font_size=18,
            )

            # Store slide index for reference
            shape._element.set("toc_target_slide", str(entry["slide_idx"]))

            y_pos += 0.7

        logger.debug("toc_added: %d entries at position %d", len(entries), insert_idx)
        return prs.slides.index(toc_slide)

    # ───────────────────────────────────────────────────────────────────────
    #  Hyperlinks
    # ───────────────────────────────────────────────────────────────────────

    def add_hyperlink(
        self,
        shape: Any,
        url: str,
        tooltip: str = "",
    ) -> bool:
        """
        Add a hyperlink to any shape.

        Args:
            shape: PPTX shape object
            url: URL to link to
            tooltip: Hover tooltip text

        Returns:
            True if successful
        """
        try:
            from pptx.enum.action import PP_ACTION

            # Validate URL
            if not self._is_valid_url(url):
                logger.warning("invalid_url: %s", url)
                return False

            # Get or create click action
            click_action = shape.click_action
            click_action.hyperlink.address = url

            if tooltip:
                # Store tooltip in alt text if possible
                try:
                    shape._element.set("tooltip", tooltip)
                except Exception:
                    pass

            return True

        except Exception as exc:
            logger.warning("hyperlink_add_failed: %s", str(exc)[:200])
            return False

    def add_source_citations(
        self,
        slide: Any,
        sources: List[Dict[str, str]],
        position: Tuple[float, float] = (0.5, 7.0),
    ) -> None:
        """
        Add source citations with clickable links.

        Args:
            slide: PPTX slide object
            sources: List of {"text": str, "url": str} dicts
            position: (left, top) position in inches
        """
        from pptx.util import Inches, Pt

        if not sources:
            return

        # Add "Sources:" label
        left, top = position
        label = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(12), Inches(0.3)
        )
        tf = label.text_frame
        tf.text = "Sources:"
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb("#64748B")

        # Add each source
        for i, source in enumerate(sources[:3]):  # Max 3 sources
            y_pos = top + 0.3 + (i * 0.25)

            source_box = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(y_pos), Inches(11), Inches(0.25)
            )
            tf = source_box.text_frame
            tf.text = f"• {source.get('text', '')}"
            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = self._hex_to_rgb("#94A3B8")

            # Make clickable if URL provided
            url = source.get("url", "")
            if url and self._is_valid_url(url):
                self.add_hyperlink(source_box, url, source.get("text", ""))

    # ───────────────────────────────────────────────────────────────────────
    #  QR Codes
    # ───────────────────────────────────────────────────────────────────────

    def generate_qr(
        self,
        data: str,
        size: int = 200,
        error_correction: str = "M",
        fill_color: str = "#000000",
        back_color: str = "#FFFFFF",
    ) -> Optional[bytes]:
        """
        Generate QR code as PNG bytes.

        Args:
            data: Data to encode (URL, contact info, etc.)
            size: Size in pixels
            error_correction: L, M, Q, or H
            fill_color: QR color hex
            back_color: Background color hex

        Returns:
            PNG bytes or None on failure
        """
        try:
            import qrcode

            # Create QR code
            qr = qrcode.QRCode(
                version=None,
                error_correction=self.QR_ERROR_CORRECTION.get(error_correction, 0),
                box_size=10,
                border=2,
            )
            qr.add_data(data)
            qr.make(fit=True)

            # Generate image
            fill_rgb = self._hex_to_rgb(fill_color)
            back_rgb = self._hex_to_rgb(back_color)

            img = qr.make_image(
                fill_color=fill_rgb,
                back_color=back_rgb,
            )

            # Resize if needed
            if size != 200:
                img = img.resize((size, size))

            # Convert to bytes
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()

        except ImportError:
            logger.warning("qrcode_module_not_installed")
            return None
        except Exception as exc:
            logger.warning("qr_generation_failed: %s", str(exc)[:200])
            return None

    def add_qr_to_slide(
        self,
        slide: Any,
        data: str,
        position: Tuple[float, float] = (11.0, 6.0),
        size: float = 1.0,
        caption: str = "",
    ) -> bool:
        """
        Generate and add QR code to a slide.

        Args:
            slide: PPTX slide object
            data: Data to encode
            position: (left, top) in inches
            size: Size in inches
            caption: Optional caption text

        Returns:
            True if successful
        """
        from pptx.util import Inches

        qr_bytes = self.generate_qr(data, size=int(size * 150))
        if not qr_bytes:
            return False

        try:
            # Add image to slide
            left, top = position
            pic = slide.shapes.add_picture(
                io.BytesIO(qr_bytes),
                Inches(left),
                Inches(top),
                Inches(size),
                Inches(size),
            )

            # Add caption if provided
            if caption:
                from pptx.util import Pt

                caption_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(top + size + 0.1),
                    Inches(size),
                    Inches(0.4),
                )
                tf = caption_box.text_frame
                tf.text = caption
                p = tf.paragraphs[0]
                p.font.size = Pt(9)
                p.font.color.rgb = self._hex_to_rgb("#64748B")
                p.alignment = 1  # Center

            return True

        except Exception as exc:
            logger.warning("qr_add_to_slide_failed: %s", str(exc)[:200])
            return False

    # ───────────────────────────────────────────────────────────────────────
    #  Navigation
    # ───────────────────────────────────────────────────────────────────────

    def add_navigation(
        self,
        prs: Any,
        style: str = "corner_buttons",
    ) -> None:
        """
        Add navigation buttons to all slides.

        Args:
            prs: Presentation object
            style: "corner_buttons", "bottom_bar", or "minimal"
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides):
            if style == "corner_buttons":
                # Previous button (except first slide)
                if idx > 0:
                    prev_btn = slide.shapes.add_shape(
                        MSO_SHAPE.CHEVRON,
                        Inches(0.5), Inches(6.8),
                        Inches(0.4), Inches(0.4),
                    )
                    prev_btn._element.set("nav_action", "prev")
                    prev_btn.fill.solid()
                    prev_btn.fill.fore_color.rgb = self._hex_to_rgb("#CBD5E1")

                # Next button (except last slide)
                if idx < total_slides - 1:
                    next_btn = slide.shapes.add_shape(
                        MSO_SHAPE.CHEVRON,
                        Inches(12.4), Inches(6.8),
                        Inches(0.4), Inches(0.4),
                    )
                    next_btn._element.set("nav_action", "next")
                    # Rotate for next direction
                    from pptx.util import Emu
                    next_btn.rotation = 180
                    next_btn.fill.solid()
                    next_btn.fill.fore_color.rgb = self._hex_to_rgb("#CBD5E1")

            elif style == "bottom_bar":
                # Progress indicator
                progress_width = (idx + 1) / total_slides * 12
                progress_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.5), Inches(7.3),
                    Inches(progress_width), Inches(0.05),
                )
                progress_bar.fill.solid()
                progress_bar.fill.fore_color.rgb = self._hex_to_rgb("#2563EB")

    # ───────────────────────────────────────────────────────────────────────
    #  Action Buttons
    # ───────────────────────────────────────────────────────────────────────

    def add_action_buttons(
        self,
        slide: Any,
        buttons: List[ActionButton],
    ) -> None:
        """
        Add action buttons to a slide.

        Args:
            slide: PPTX slide object
            buttons: List of ActionButton definitions
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        for btn in buttons:
            left, top = btn.position
            width, height = btn.size

            # Create button shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(height),
            )

            # Style button
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(btn.color)
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.text = btn.label
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
            p.alignment = 1  # Center

            # Store action
            action_data = f"{btn.action}:{btn.target}"
            shape._element.set("action", action_data)

            # Add hyperlink if URL
            if btn.action == "url":
                self.add_hyperlink(shape, btn.target)
            elif btn.action == "email":
                mailto = f"mailto:{btn.target}"
                self.add_hyperlink(shape, mailto)

    def create_contact_buttons(
        self,
        slide: Any,
        email: str = "",
        website: str = "",
        phone: str = "",
        position: Tuple[float, float] = (4.0, 6.0),
    ) -> None:
        """
        Add standard contact action buttons.

        Args:
            slide: PPTX slide object
            email: Contact email
            website: Website URL
            phone: Phone number
            position: Starting position for buttons
        """
        buttons = []
        x_pos = position[0]

        if website:
            buttons.append(ActionButton(
                label="Visit Website",
                action="url",
                target=website,
                color="#2563EB",
                position=(x_pos, position[1]),
            ))
            x_pos += 2.5

        if email:
            buttons.append(ActionButton(
                label="Contact Us",
                action="email",
                target=email,
                color="#16A34A",
                position=(x_pos, position[1]),
            ))
            x_pos += 2.5

        if phone:
            buttons.append(ActionButton(
                label="Call Now",
                action="url",
                target=f"tel:{phone}",
                color="#DC2626",
                position=(x_pos, position[1]),
            ))

        if buttons:
            self.add_action_buttons(slide, buttons)

    # ───────────────────────────────────────────────────────────────────────
    #  Video Embeds
    # ───────────────────────────────────────────────────────────────────────

    def add_video_placeholder(
        self,
        slide: Any,
        video_url: str,
        position: Tuple[float, float, float, float] = (1.0, 2.0, 11.0, 5.0),
        thumbnail_url: Optional[str] = None,
    ) -> bool:
        """
        Add a video placeholder with link to external video.

        Args:
            slide: PPTX slide object
            video_url: YouTube, Vimeo, or direct video URL
            position: (left, top, width, height) in inches
            thumbnail_url: Optional custom thumbnail

        Returns:
            True if successful
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        try:
            left, top, width, height = position

            # Create placeholder shape
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(height),
            )

            # Style as video placeholder
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self._hex_to_rgb("#1E293B")
            placeholder.line.color.rgb = self._hex_to_rgb("#475569")

            # Add play icon (triangle)
            icon_size = min(width, height) * 0.2
            icon_x = left + width / 2 - icon_size / 2
            icon_y = top + height / 2 - icon_size / 2

            play_icon = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(icon_x), Inches(icon_y),
                Inches(icon_size), Inches(icon_size),
            )
            play_icon.fill.solid()
            play_icon.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF")
            play_icon.line.fill.background()

            # Add "Watch Video" text
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + height - 0.8),
                Inches(width), Inches(0.5),
            )
            tf = text_box.text_frame
            tf.text = "▶ Click to Watch Video"
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
            p.alignment = 1  # Center

            # Add hyperlink
            self.add_hyperlink(placeholder, video_url, "Watch Video")
            self.add_hyperlink(play_icon, video_url, "Watch Video")

            return True

        except Exception as exc:
            logger.warning("video_placeholder_add_failed: %s", str(exc)[:200])
            return False

    # ───────────────────────────────────────────────────────────────────────
    #  Utility Methods
    # ───────────────────────────────────────────────────────────────────────

    def _is_valid_url(self, url: str) -> bool:
        """Validate URL format."""
        if not url:
            return False

        url = url.strip()

        # Basic URL validation
        patterns = [
            r'^https?://',  # HTTP(S)
            r'^mailto:',     # Email
            r'^tel:',        # Phone
            r'^file:',       # File
        ]

        return any(re.match(pattern, url, re.IGNORECASE) for pattern in patterns)

    def _add_text_box(
        self,
        slide: Any,
        text: str,
        left: Any,
        top: Any,
        width: Any,
        height: Any,
        font_size: int = 18,
        bold: bool = False,
        color: str = "#0F172A",
    ) -> Any:
        """Helper to add a text box."""
        from pptx.util import Pt

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._hex_to_rgb(color)

        return shape

    def _hex_to_rgb(self, value: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        h = value.strip().lstrip("#")
        if len(h) != 6:
            return (37, 99, 235)  # Default blue
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Convenience functions

def add_toc_to_presentation(prs: Any, slides: List[SlideSpec]) -> int:
    """One-shot TOC addition."""
    builder = InteractiveBuilder()
    return builder.add_toc(prs, slides)


def generate_qr_code(data: str, size: int = 200) -> Optional[bytes]:
    """One-shot QR generation."""
    builder = InteractiveBuilder()
    return builder.generate_qr(data, size)


def add_contact_buttons(
    slide: Any,
    email: str = "",
    website: str = "",
    phone: str = "",
) -> None:
    """One-shot contact buttons."""
    builder = InteractiveBuilder()
    builder.create_contact_buttons(slide, email, website, phone)
