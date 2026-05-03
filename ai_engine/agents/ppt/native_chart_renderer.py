"""
NativeChartRenderer — python-pptx native Chart objects (editable in PowerPoint).

Unlike ChartRenderer (matplotlib PNGs), this embeds actual data tables that
users can edit: right-click chart → Edit Data.

Supported kinds:
    bar, column, line, pie, scatter, bubble

Design notes:
- Chart data is embedded as CategoryChartData / XyChartData / BubbleChartData
- Styling follows the deck theme (accent color, fonts)
- Falls back to PNG if native rendering fails
- All series must have same length for native charts (enforced)
"""
from __future__ import annotations

import logging
from typing import Any, Dict, List, Optional, Union

from pptx.chart.data import (
    BubbleChartData,
    CategoryChartData,
    ChartData,
    XyChartData,
)
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from ai_engine.agents.ppt.schemas import ChartKind, ChartSpec

logger = logging.getLogger(__name__)

# Map our ChartKind to python-pptx XL_CHART_TYPE
_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble": XL_CHART_TYPE.BUBBLE,
    # Stacked variants
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    # Area
    "area": XL_CHART_TYPE.AREA,
    "stacked_area": XL_CHART_TYPE.AREA_STACKED,
}


def _hex_to_rgb(value: str) -> tuple:
    """Convert #RRGGBB to (R, G, B) ints."""
    h = value.strip().lstrip("#")
    if len(h) != 6:
        return (37, 99, 235)  # default blue (#2563EB)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _palette(accent_hex: str, n: int) -> List[tuple]:
    """Build n-color palette, accent first then complements."""
    accent = _hex_to_rgb(accent_hex)
    extras = [
        (239, 68, 68),   # red
        (34, 197, 94),   # green
        (245, 158, 11),  # amber
        (139, 92, 246),  # purple
        (6, 182, 212),   # cyan
        (236, 72, 153),  # pink
        (107, 114, 128), # gray
    ]
    out = [accent]
    for c in extras:
        if len(out) >= n:
            break
        out.append(c)
    while len(out) < n:
        out.append(extras[len(out) % len(extras)])
    return out[:n]


class NativeChartRenderer:
    """Render ChartSpec as native editable PowerPoint charts."""

    def __init__(self, dpi: int = 150) -> None:
        self.dpi = dpi

    def can_render(self, spec: ChartSpec) -> bool:
        """Check if we can render this spec natively."""
        kind = (spec.kind or "").lower().strip()
        if kind not in _CHART_TYPE_MAP:
            return False
        # Native charts need at least one series
        if not spec.series:
            return False
        # All series must have data
        for s in spec.series:
            if not s.get("data"):
                return False
        return True

    def render_to_slide(
        self,
        slide,
        spec: ChartSpec,
        *,
        accent_hex: str = "#2563EB",
        left: Any = None,
        top: Any = None,
        width: Any = None,
        height: Any = None,
    ) -> Optional[Any]:
        """
        Add a native chart to the slide.

        Args:
            slide: pptx Slide object
            spec: ChartSpec with data
            accent_hex: theme accent color
            left, top, width, height: Inches or Pt positions

        Returns:
            The chart shape or None on failure
        """
        kind = (spec.kind or "").lower().strip()
        chart_type = _CHART_TYPE_MAP.get(kind)
        if chart_type is None:
            logger.debug("native_chart_unsupported_kind: %s", kind)
            return None

        # Default positioning
        if left is None:
            left = Inches(1.0)
        if top is None:
            top = Inches(1.9)
        if width is None:
            width = Inches(11.3)
        if height is None:
            height = Inches(4.6)

        try:
            chart_data = self._build_chart_data(spec, accent_hex)
            if chart_data is None:
                return None

            chart = slide.shapes.add_chart(
                chart_type,
                left, top, width, height,
                chart_data,
            ).chart

            # Apply styling
            self._style_chart(chart, spec, accent_hex)
            return chart

        except Exception as exc:
            logger.warning("native_chart_render_failed: %s", str(exc)[:200])
            return None

    def _build_chart_data(
        self, spec: ChartSpec, accent_hex: str
    ) -> Optional[Union[ChartData, CategoryChartData, XyChartData, BubbleChartData]]:
        """Build appropriate ChartData subclass from spec."""
        kind = (spec.kind or "").lower().strip()

        if kind in ("scatter",):
            return self._build_xy_chart_data(spec)
        if kind in ("bubble",):
            return self._build_bubble_chart_data(spec)
        # Default: category-based (bar, column, line, pie, area)
        return self._build_category_chart_data(spec)

    def _build_category_chart_data(self, spec: ChartSpec) -> Optional[CategoryChartData]:
        """Build CategoryChartData for bar/column/line/pie/area."""
        if not spec.series:
            return None

        chart_data = CategoryChartData()

        # Set categories (x-axis labels)
        categories = spec.categories or []
        if not categories:
            # Infer from first series length
            first_data = spec.series[0].get("data") or []
            categories = [f"Item {i+1}" for i in range(len(first_data))]
        chart_data.categories = categories

        # Add series
        for i, s in enumerate(spec.series):
            name = s.get("name") or f"Series {i+1}"
            raw_data = s.get("data") or []
            # Extract numeric values
            values = []
            for d in raw_data:
                if isinstance(d, (int, float)):
                    values.append(float(d))
                elif isinstance(d, (list, tuple)) and len(d) >= 2:
                    values.append(float(d[1]))  # (x, y) -> y
                else:
                    values.append(0.0)
            # Pad or truncate to match categories
            while len(values) < len(categories):
                values.append(0.0)
            values = values[:len(categories)]

            # Add series
            chart_data.add_series(name, values)

        return chart_data

    def _build_xy_chart_data(self, spec: ChartSpec) -> Optional[XyChartData]:
        """Build XyChartData for scatter plots."""
        if not spec.series:
            return None

        chart_data = XyChartData()

        for i, s in enumerate(spec.series):
            name = s.get("name") or f"Series {i+1}"
            raw_data = s.get("data") or []
            series = chart_data.add_series(name)

            for d in raw_data:
                if isinstance(d, (list, tuple)) and len(d) >= 2:
                    x, y = float(d[0]), float(d[1])
                    series.add_data_point(x, y)
                elif isinstance(d, (int, float)):
                    # Single value: use index as x
                    series.add_data_point(float(i), float(d))

        return chart_data

    def _build_bubble_chart_data(self, spec: ChartSpec) -> Optional[BubbleChartData]:
        """Build BubbleChartData for bubble charts."""
        if not spec.series:
            return None

        chart_data = BubbleChartData()

        for i, s in enumerate(spec.series):
            name = s.get("name") or f"Series {i+1}"
            raw_data = s.get("data") or []
            series = chart_data.add_series(name)

            for d in raw_data:
                if isinstance(d, (list, tuple)) and len(d) >= 3:
                    x, y, size = float(d[0]), float(d[1]), float(d[2])
                    series.add_data_point(x, y, size)
                elif isinstance(d, (list, tuple)) and len(d) >= 2:
                    x, y = float(d[0]), float(d[1])
                    series.add_data_point(x, y, 20.0)  # default size
                elif isinstance(d, (int, float)):
                    series.add_data_point(float(i), float(d), 20.0)

        return chart_data

    def _style_chart(self, chart, spec: ChartSpec, accent_hex: str) -> None:
        """Apply consistent styling to the chart."""
        # Title
        if spec.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = spec.title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True
        else:
            chart.has_title = False

        # Legend (show if multiple series)
        if len(spec.series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False

        # Category axis label (x-axis)
        if spec.x_label and hasattr(chart, "category_axis"):
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = spec.x_label
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

        # Value axis label (y-axis)
        if spec.y_label and hasattr(chart, "value_axis"):
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = spec.y_label
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

        # Style the plot area
        chart.plot_area.format.fill.background()  # transparent

        # Gridlines (subtle)
        if hasattr(chart, "value_axis"):
            chart.value_axis.major_gridlines.format.line.fill.background()


# Convenience function for direct use
def render_native_chart(
    slide,
    spec: ChartSpec,
    *,
    accent_hex: str = "#2563EB",
    left=None,
    top=None,
    width=None,
    height=None,
) -> Optional[Any]:
    """One-shot native chart render."""
    renderer = NativeChartRenderer()
    return renderer.render_to_slide(
        slide, spec,
        accent_hex=accent_hex,
        left=left, top=top, width=width, height=height,
    )
